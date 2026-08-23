"""PowerPoint → Hanns JSON importer (v2 — high-fidelity).

Imports PPTX (and legacy PPT via LibreOffice conversion when available) into
the existing Hanns Deck/Slide JSON model. Images are stored as media files and
only URLs live inside slide JSON so autosave stays small.

What v2 imports beyond the original importer:

  • Slide backgrounds  — solid fills, real CSS gradients built from the PPTX
    gradient stops + angle, and full-bleed background pictures (slide →
    layout → master inheritance).
  • Text               — bullet prefixes, per-shape line spacing, alignment,
    fonts, bold/italic/colour, and text-box fills; a text box whose whole
    content is one hyperlink becomes a Hanns link button.
  • Tables             — real PPTX tables become editable Hanns `table`
    elements with their cell data.
  • Charts             — real categories + values are extracted (first
    series, or all series → grouped bars), the chart type is mapped to the
    matching Hanns `chartKind`, and series colours seed the palette.
  • Lines              — true stroke width and dash style.
  • Speaker notes      — imported per slide.
  • Entrance animations— the slide's <p:timing> tree is parsed and PowerPoint
    entrance presets (Fly In, Fade, Zoom, Wipe, Swivel, Wheel, …) are mapped
    onto the matching Hanns entrance (`anim`), with their delay (`animDelay`)
    and duration (`animDur`) preserved.
  • Z-order            — element order follows PPTX document order exactly.

v3 fidelity fixes (import now matches the source PowerPoint):

  • Font sizing        — sizes are scaled to the slide's real width so a
    29 pt title on a 13.33-inch slide lands at ~29 px in the 960×540 design
    space (the old fixed ×1.333 was only correct for 10-inch slides and
    made every 16:9 deck ~33 % oversized, overflowing and overlapping).
  • normAutofit        — PowerPoint's "shrink text on overflow" fontScale
    and lnSpcReduction are honored, so autofitted boxes stay inside their
    frames.
  • No-fill shapes     — an AUTO_SHAPE whose fill is "no fill"
    (MSO_FILL.BACKGROUND) no longer receives the beige placeholder fill;
    invisible text-container shapes are skipped entirely instead of
    painting stray boxes behind every text run.
  • Line geometry      — shapes whose preset geometry is a line/connector
    (prst="line", straightConnector1, …) import as real Hanns line
    elements with their stroke colour, instead of 4 px-tall beige rects.
  • Vertical anchor    — the shape's bodyPr anchor (top/middle/bottom) is
    carried on text elements as ``valign``.
  • Dash detection     — no longer mutates the source XML while reading.
"""

from __future__ import annotations

import os
import re
import subprocess
import tempfile
import uuid
from pathlib import Path

from django.core.files.base import ContentFile
from django.core.files.storage import default_storage

try:  # python-pptx is optional until the import feature is used.
    from pptx import Presentation
    from pptx.enum.dml import MSO_FILL, MSO_COLOR_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
except Exception:  # pragma: no cover - handled at runtime
    Presentation = None
    MSO_FILL = None
    MSO_COLOR_TYPE = None
    MSO_SHAPE_TYPE = None
    PP_ALIGN = None
    qn = None

from .models import Slide

DESIGN_W = 960
DESIGN_H = 540


# ─────────────────────────── small utilities ───────────────────────────

def _uid() -> str:
    return "ppt_" + uuid.uuid4().hex[:12]


def _safe_name(name: str) -> str:
    base = os.path.basename(name or "presentation.pptx")
    return base.replace("/", "_").replace("\\", "_")[:160]


def _write_upload_to_temp(uploaded_file, path: str) -> None:
    with open(path, "wb") as fh:
        for chunk in uploaded_file.chunks():
            fh.write(chunk)


def _convert_legacy_ppt_to_pptx(source_path: str, outdir: str) -> str:
    """Use LibreOffice to convert old .ppt to .pptx when available."""
    cmd = [
        "soffice", "--headless", "--convert-to", "pptx",
        "--outdir", outdir, source_path,
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=60)
    if proc.returncode != 0:
        raise ValueError(
            "Could not convert this legacy .ppt file. Install LibreOffice in the web container "
            "or upload a .pptx version."
        )
    candidates = list(Path(outdir).glob("*.pptx"))
    if not candidates:
        raise ValueError("PowerPoint conversion completed but no .pptx file was produced.")
    return str(candidates[0])


def _load_presentation(uploaded_file):
    if Presentation is None:
        raise ValueError("python-pptx is not installed. Add python-pptx to requirements.txt and rebuild.")

    original_name = _safe_name(getattr(uploaded_file, "name", "presentation.pptx"))
    ext = Path(original_name).suffix.lower()
    if ext not in {".pptx", ".ppt", ".ppsx"}:
        raise ValueError("Only .ppt, .pptx, and .ppsx PowerPoint files can be imported.")

    tmpdir = tempfile.TemporaryDirectory()
    source_path = os.path.join(tmpdir.name, "source" + ext)
    _write_upload_to_temp(uploaded_file, source_path)
    parse_path = source_path
    if ext == ".ppt":
        parse_path = _convert_legacy_ppt_to_pptx(source_path, tmpdir.name)
    return Presentation(parse_path), original_name, tmpdir


def _emu(value, total, design):
    try:
        return round((int(value) / int(total)) * design, 2) if total else 0
    except Exception:
        return 0


def _box(shape, prs):
    return {
        "x": _emu(getattr(shape, "left", 0), prs.slide_width, DESIGN_W),
        "y": _emu(getattr(shape, "top", 0), prs.slide_height, DESIGN_H),
        "w": max(4, _emu(getattr(shape, "width", 0), prs.slide_width, DESIGN_W)),
        "h": max(4, _emu(getattr(shape, "height", 0), prs.slide_height, DESIGN_H)),
        "rot": float(getattr(shape, "rotation", 0) or 0),
    }


def _media_url(saved_path: str) -> str:
    """URL for a freshly saved media file, as the BROWSER should see it.

    Deliberately NOT ``request.build_absolute_uri(...)``. Behind a reverse
    proxy / CDN, build_absolute_uri regularly produces the wrong scheme
    (http:// on an https:// page → blocked as mixed content) or an internal
    hostname the browser cannot reach. Either way the slide element renders
    as an EMPTY FRAME, because CSS background-image fails silently.

    A root-relative "/media/…" URL always resolves against the page origin,
    so it works in dev, behind nginx, and behind Cloudflare alike. Remote
    storages (S3, GCS) already return a fully-qualified URL from
    ``default_storage.url()``; those are passed through untouched.
    """
    url = default_storage.url(saved_path)
    if url.startswith(("http://", "https://", "//", "data:")):
        return url
    return url if url.startswith("/") else "/" + url


def _shape_image(shape):
    """Return the shape's embedded image, for ANY shape that carries one.

    ``shape_type == PICTURE`` is not the only way a picture reaches a slide.
    Pictures dropped into a layout's picture/content placeholder keep
    ``shape_type == PLACEHOLDER`` (python-pptx: PicturePlaceholder), and
    linked pictures report LINKED_PICTURE. Both were previously skipped by
    the importer, so template-based decks lost every image.
    """
    try:
        image = shape.image
    except (AttributeError, ValueError, KeyError, TypeError):
        return None
    return image if getattr(image, "blob", None) else None


def _pt_scale(prs) -> float:
    """px-per-pt when mapping this deck into the 960×540 design space.

    A PPTX point is 12700 EMU. Scaling by the slide's real width means a
    16:9 13.33-inch slide (960 pt wide) maps 1 pt → 1 px, while an old 4:3
    10-inch slide (720 pt wide) maps 1 pt → 1.333 px — exactly the fixed
    factor the previous importer applied to everything, which oversized
    all widescreen decks by ~33 %.
    """
    try:
        slide_w_pt = int(prs.slide_width) / 12700.0
        return (DESIGN_W / slide_w_pt) if slide_w_pt > 0 else 1.333
    except Exception:
        return 1.333


# ─────────────────────────── colour helpers ───────────────────────────

def _rgb_to_hex(rgb) -> str | None:
    if rgb is None:
        return None
    value = str(rgb).strip()
    if len(value) == 6:
        return "#" + value.upper()
    return None


def _color_to_hex(color, default=None):
    if color is None:
        return default
    try:
        if getattr(color, "type", None) == MSO_COLOR_TYPE.RGB:
            return _rgb_to_hex(color.rgb) or default
        # Theme colours do not resolve reliably without the PowerPoint theme;
        # return the default instead of throwing.
        return default
    except Exception:
        return default


def _srgb_from_xml(clr_el) -> str | None:
    """Extract a hex colour from an <a:srgbClr>/<a:sysClr> child of clr_el."""
    if clr_el is None:
        return None
    srgb = clr_el.find(qn("a:srgbClr"))
    if srgb is not None and srgb.get("val"):
        return "#" + srgb.get("val").upper()
    sysc = clr_el.find(qn("a:sysClr"))
    if sysc is not None and sysc.get("lastClr"):
        return "#" + sysc.get("lastClr").upper()
    return None


def _fill_kind(shape):
    """Classify a shape's fill without guessing.

    Returns (kind, hex) where kind is:
      "none"  — explicitly no fill (MSO_FILL.BACKGROUND) or none set
      "solid" — solid fill; hex is the colour when it resolved to RGB
      "other" — gradient / pattern / picture fill
    """
    try:
        fill = shape.fill
        ftype = getattr(fill, "type", None)
        if ftype is None or (MSO_FILL is not None and ftype == MSO_FILL.BACKGROUND):
            return "none", None
        if MSO_FILL is not None and ftype == MSO_FILL.SOLID:
            return "solid", _color_to_hex(fill.fore_color, None)
        return "other", None
    except Exception:
        return "none", None


def _fill_color(shape, default="none"):
    """A shape with NO fill must import as "none" — never as a beige box.

    The default only applies when the shape genuinely has a fill Hanns
    cannot represent (theme-coloured solid, gradient, pattern).
    """
    kind, hexv = _fill_kind(shape)
    if kind == "solid":
        return hexv or default
    if kind == "other":
        return default
    return "none"


def _fill_or(shape, fallback):
    """Solid fill colour, or ``fallback`` when there is none/unresolvable."""
    kind, hexv = _fill_kind(shape)
    return hexv if (kind == "solid" and hexv) else fallback


def _stroke_color(shape, default="none"):
    try:
        line = shape.line
        if line is None:
            return default
        # No outline (line fill is BACKGROUND / not set) → no stroke.
        lft = getattr(getattr(line, "fill", None), "type", None)
        if MSO_FILL is not None and lft == MSO_FILL.BACKGROUND:
            return "none"
        if line.color is not None:
            return _color_to_hex(line.color, default)
    except Exception:
        pass
    return default


def _stroke_width(shape, scale=1.333):
    try:
        if shape.line.width:
            return max(1, round(float(shape.line.width.pt) * scale, 1))
        return 0
    except Exception:
        return 0


def _line_is_dashed(shape) -> bool:
    # Read-only lookup — the old _get_or_add_ln() call MUTATED the source
    # XML while merely checking for a dash style.
    try:
        ln = shape._element.find(".//" + qn("a:ln"))
        if ln is None:
            return False
        dash = ln.find(qn("a:prstDash"))
        return dash is not None and (dash.get("val") or "solid") != "solid"
    except Exception:
        return False


# Preset geometries that are really lines/connectors. PowerPoint files
# frequently store these as AUTO_SHAPEs (not MSO_SHAPE_TYPE.LINE), which the
# old importer turned into 4 px-tall filled rectangles — losing every
# divider line in the deck.
_LINE_PRSTS = {
    "line", "lineInv",
    "straightConnector1",
    "bentConnector2", "bentConnector3", "bentConnector4", "bentConnector5",
    "curvedConnector2", "curvedConnector3", "curvedConnector4", "curvedConnector5",
}


def _is_line_geometry(shape) -> bool:
    try:
        pg = shape._element.find(".//" + qn("a:prstGeom"))
        return pg is not None and (pg.get("prst") or "") in _LINE_PRSTS
    except Exception:
        return False


# ─────────────────────────── background import ───────────────────────────

def _gradient_css_from_xml(grad_el) -> str | None:
    """Build a CSS linear-gradient from an <a:gradFill> element."""
    try:
        stops = []
        gs_lst = grad_el.find(qn("a:gsLst"))
        if gs_lst is None:
            return None
        for gs in gs_lst.findall(qn("a:gs")):
            hexv = _srgb_from_xml(gs)
            if not hexv:
                continue
            pos = int(gs.get("pos") or 0) / 1000.0  # 0–100 (%)
            stops.append((pos, hexv))
        if len(stops) < 2:
            return None
        stops.sort(key=lambda s: s[0])
        # PPTX angle is clockwise from 3 o'clock in 60000ths of a degree;
        # CSS gradients are clockwise from 12 o'clock.
        angle = 135
        lin = grad_el.find(qn("a:lin"))
        if lin is not None and lin.get("ang"):
            angle = (int(lin.get("ang")) / 60000.0 + 90) % 360
        parts = ",".join(f"{c} {p:.0f}%" for p, c in stops)
        return f"linear-gradient({angle:.0f}deg,{parts})"
    except Exception:
        return None


def _bg_picture_url(request, deck, bg_pr, part) -> str | None:
    """Resolve an <a:blipFill> background picture to a stored media URL."""
    try:
        blip_fill = bg_pr.find(qn("a:blipFill"))
        if blip_fill is None:
            return None
        blip = blip_fill.find(qn("a:blip"))
        if blip is None:
            return None
        r_id = blip.get(qn("r:embed"))
        if not r_id:
            return None
        image_part = part.related_parts[r_id]
        ext = (image_part.partname.ext or "png").lower().lstrip(".")
        if ext == "jpeg":
            ext = "jpg"
        rel_path = f"hanns/decks/{deck.code}/imports/bg_{uuid.uuid4().hex}.{ext}"
        saved = default_storage.save(rel_path, ContentFile(image_part.blob))
        return _media_url(saved)
    except Exception:
        return None


def _slide_background(request, deck, slide) -> tuple[str, str | None]:
    """Return (bg CSS, bgSize) for a slide, walking slide → layout → master."""
    parts = [slide]
    try:
        parts.append(slide.slide_layout)
        parts.append(slide.slide_layout.slide_master)
    except Exception:
        pass

    for holder in parts:
        try:
            root = holder._element
            bg = root.find(qn("p:cSld") + "/" + qn("p:bg"))
            if bg is None:
                continue
            bg_pr = bg.find(qn("p:bgPr"))
            if bg_pr is None:
                continue
            # solid fill
            solid = bg_pr.find(qn("a:solidFill"))
            if solid is not None:
                hexv = _srgb_from_xml(solid)
                if hexv:
                    return hexv, None
            # gradient fill
            grad = bg_pr.find(qn("a:gradFill"))
            if grad is not None:
                css = _gradient_css_from_xml(grad)
                if css:
                    return css, None
            # picture fill
            url = _bg_picture_url(request, deck, bg_pr, holder.part)
            if url:
                return f'url("{url}") center/cover no-repeat', "cover"
        except Exception:
            continue
    return "#f6f1e7", None


# ─────────────────────────── animation import ───────────────────────────

# presetSubtype for directional effects is a bitfield: 1=top 2=right 4=bottom 8=left.
def _direction_anim(subtype: int) -> str:
    if subtype & 4:
        return "rise"    # flies in from the bottom
    if subtype & 1:
        return "drop"    # from the top
    if subtype & 8:
        return "left"    # from the left
    if subtype & 2:
        return "right"   # from the right
    return "rise"


def _preset_to_anim(preset_id: int, subtype: int) -> str:
    if preset_id in (1, 9, 10):            # Appear / Dissolve / Fade
        return "fade"
    if preset_id in (2, 7, 12):            # Fly In / Crawl In / Peek In
        return _direction_anim(subtype)
    if preset_id == 22:                    # Wipe
        return "revealUp" if subtype & (1 | 4) else "reveal"
    if preset_id in (3, 5, 14, 18):        # Blinds / Checkerboard / Random bars / Strips
        return "reveal"
    if preset_id in (4, 6, 8, 13, 16, 20, 23):  # Box/Circle/Diamond/Plus/Split/Wedge/Zoom
        return "zoom"
    if preset_id == 17:                    # Stretch
        return "pop"
    if preset_id == 19:                    # Swivel
        return "flipx"
    if preset_id in (15, 21):              # Spiral / Wheel
        return "spin"
    return "fade"


def _entrance_animations(slide) -> dict[int, dict]:
    """Parse the slide's <p:timing> tree.

    Returns {shape_id: {"anim": str, "animDelay": float, "animDur": float}}
    for every shape that has a PowerPoint entrance effect.
    """
    out: dict[int, dict] = {}
    try:
        timing = slide._element.find(qn("p:timing"))
        if timing is None:
            return out
        running_delay = 0.0
        for ctn in timing.iter(qn("p:cTn")):
            if ctn.get("presetClass") != "entr":
                continue
            try:
                preset_id = int(ctn.get("presetID") or 0)
            except ValueError:
                preset_id = 0
            try:
                subtype = int(ctn.get("presetSubtype") or 0)
            except ValueError:
                subtype = 0

            # Delay: the effect node's own start condition (ms).
            delay_ms = 0.0
            st = ctn.find(qn("p:stCondLst"))
            if st is not None:
                cond = st.find(qn("p:cond"))
                if cond is not None:
                    raw = cond.get("delay")
                    if raw and raw != "indefinite":
                        try:
                            delay_ms = float(raw)
                        except ValueError:
                            delay_ms = 0.0

            # Duration: longest child behaviour duration.
            dur_ms = 0.0
            for beh in ctn.iter(qn("p:cTn")):
                if beh is ctn:
                    continue
                raw = beh.get("dur")
                if raw and raw not in ("indefinite", "0"):
                    try:
                        dur_ms = max(dur_ms, float(raw))
                    except ValueError:
                        pass

            # Target shape id.
            sp_tgt = ctn.find(".//" + qn("p:spTgt"))
            if sp_tgt is None:
                continue
            try:
                spid = int(sp_tgt.get("spid"))
            except (TypeError, ValueError):
                continue

            entry = {
                "anim": _preset_to_anim(preset_id, subtype),
                "animDelay": round(min(8.0, (running_delay + delay_ms) / 1000.0), 2),
            }
            if dur_ms:
                entry["animDur"] = round(min(5.0, dur_ms / 1000.0), 2)
            out.setdefault(spid, entry)
        return out
    except Exception:
        return out


# ─────────────────────────── text import ───────────────────────────

def _paragraph_align(paragraph):
    try:
        if paragraph.alignment == PP_ALIGN.CENTER:
            return "center"
        if paragraph.alignment == PP_ALIGN.RIGHT:
            return "right"
    except Exception:
        pass
    return "left"


def _dominant_run(shape):
    """The run carrying the most characters — best proxy for the box's style."""
    best_p, best_r, best_len = None, None, -1
    try:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                n = len(r.text or "")
                if n > best_len:
                    best_p, best_r, best_len = p, r, n
        if best_p is None and shape.text_frame.paragraphs:
            best_p = shape.text_frame.paragraphs[0]
    except Exception:
        return None, None
    return best_p, best_r


def _font_size(run, scale=1.333, fallback=28):
    try:
        if run and run.font and run.font.size:
            return max(6, int(round(float(run.font.size.pt) * scale)))
    except Exception:
        pass
    return fallback


def _autofit(shape):
    """PowerPoint "shrink text on overflow" → (fontScale 0–1, lnSpcReduction 0–1).

    <a:normAutofit fontScale="62500" lnSpcReduction="20000"/> means the text
    is rendered at 62.5 % size with line spacing reduced by 20 %. Ignoring
    this made autofitted boxes overflow their frames on import.
    """
    try:
        body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        if body_pr is None:
            return 1.0, 0.0
        na = body_pr.find(qn("a:normAutofit"))
        if na is None:
            return 1.0, 0.0
        fs = na.get("fontScale")
        lr = na.get("lnSpcReduction")
        font_scale = (int(fs) / 100000.0) if fs else 1.0
        lh_reduction = (int(lr) / 100000.0) if lr else 0.0
        return max(0.1, min(1.0, font_scale)), max(0.0, min(0.8, lh_reduction))
    except Exception:
        return 1.0, 0.0


def _text_valign(shape) -> str:
    """The shape's vertical text anchor: top | middle | bottom."""
    try:
        body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        anchor = body_pr.get("anchor") if body_pr is not None else None
        return {"ctr": "middle", "b": "bottom"}.get(anchor, "top")
    except Exception:
        return "top"


def _font_name(run):
    try:
        name = run.font.name if run and run.font else None
        if name:
            return f'"{name}",sans-serif'
    except Exception:
        pass
    return '"Inter",sans-serif'


def _font_color(run, default="#16140f"):
    try:
        if run and run.font:
            return _color_to_hex(run.font.color, default)
    except Exception:
        pass
    return default


def _paragraph_has_bullet(paragraph) -> bool:
    try:
        pPr = paragraph._pPr
        if pPr is None:
            return False
        if pPr.find(qn("a:buNone")) is not None:
            return False
        if pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None:
            return True
        # Indented levels default to bulleted in most PowerPoint themes.
        return int(pPr.get("lvl") or 0) > 0
    except Exception:
        return False


def _line_spacing(shape) -> float:
    try:
        for p in shape.text_frame.paragraphs:
            ls = p.line_spacing
            if ls is None:
                continue
            if isinstance(ls, float):        # multiple of line height
                return round(max(0.8, min(3.0, ls)), 2)
            return 1.15                       # exact point spacing → sane default
    except Exception:
        pass
    return 1.15


def _shape_text(shape, with_bullets: bool = True):
    try:
        lines = []
        for p in shape.text_frame.paragraphs:
            txt = "".join(r.text or "" for r in p.runs) or p.text or ""
            if with_bullets and txt.strip() and _paragraph_has_bullet(p):
                indent = "   " * max(0, int(getattr(p, "level", 0) or 0) - 0)
                txt = f"{indent}• {txt}"
            lines.append(txt)
        return "\n".join(lines).strip("\n").strip()
    except Exception:
        return ""


def _whole_shape_hyperlink(shape) -> str | None:
    """If every run in the shape shares one hyperlink, return its address."""
    try:
        urls = set()
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if not (r.text or "").strip():
                    continue
                addr = getattr(getattr(r, "hyperlink", None), "address", None)
                if not addr:
                    return None
                urls.add(addr)
        if len(urls) == 1:
            return urls.pop()
    except Exception:
        pass
    return None


def _text_element(shape, prs, anims):
    text = _shape_text(shape)
    if not text:
        return None
    p, r = _dominant_run(shape)
    box = _box(shape, prs)
    fill = _fill_color(shape, "none")
    scale = _pt_scale(prs)
    autofit_scale, lh_reduction = _autofit(shape)

    url = _whole_shape_hyperlink(shape)
    if url:
        # A pure-link text box becomes a proper Hanns link button.
        return {
            "id": _uid(),
            "type": "link",
            **box,
            "anim": "rise",
            "animDelay": 0,
            "url": url,
            "label": text.split("\n")[0][:80],
            "description": url,
            "linkStyle": "button",
            "accent": _font_color(r, "#2563eb"),
            "textColor": "#ffffff",
            "bg": _font_color(r, "#2563eb"),
            "radius": 18,
            **_anim_for(shape, anims),
        }

    return {
        "id": _uid(),
        "type": "text",
        **box,
        "anim": "fade",
        "animDelay": 0,
        "text": text,
        "font": _font_name(r),
        "size": _font_size(
            r,
            scale=scale * autofit_scale,
            fallback=max(14, min(44, int(box["h"] / 2) if box["h"] else 28)),
        ),
        "weight": 700 if getattr(getattr(r, "font", None), "bold", False) else 500,
        "italic": bool(getattr(getattr(r, "font", None), "italic", False)),
        "color": _font_color(r, "#16140f"),
        "align": _paragraph_align(p),
        "valign": _text_valign(shape),
        "lh": max(0.9, round(_line_spacing(shape) * (1.0 - lh_reduction), 2)),
        "ls": 0,
        "fill": fill,
        **_anim_for(shape, anims),
    }


# ─────────────────────────── media / shapes / tables / charts ───────────────

def _anim_for(shape, anims: dict[int, dict]) -> dict:
    try:
        entry = anims.get(int(shape.shape_id))
        return dict(entry) if entry else {}
    except Exception:
        return {}


def _image_element(request, deck, shape, prs, anims):
    try:
        image = _shape_image(shape)
        if image is None:
            return None
        ext = (image.ext or "png").lower().lstrip(".")
        if ext == "jpeg":
            ext = "jpg"
        rel_path = f"hanns/decks/{deck.code}/imports/{uuid.uuid4().hex}.{ext}"
        saved = default_storage.save(rel_path, ContentFile(image.blob))
        url = _media_url(saved)
        return {
            "id": _uid(),
            "type": "image",
            **_box(shape, prs),
            "anim": "fade",
            "animDelay": 0,
            "src": url,
            "fit": "contain",
            "radius": 0,
            "alt": "Imported PowerPoint image",
            **_anim_for(shape, anims),
        }
    except Exception:
        return None


def _shape_element(shape, prs, anims):
    box = _box(shape, prs)
    # "none" for genuinely unfilled shapes; the neutral default only kicks
    # in for fills Hanns cannot represent (gradient/pattern/theme solid).
    fill = _fill_color(shape, "#ece4d4")
    stroke = _stroke_color(shape, "none")
    stroke_w = _stroke_width(shape, _pt_scale(prs))
    typ = "rect"
    radius = 8
    try:
        auto_name = str(getattr(shape, "auto_shape_type", "")).upper()
        if "OVAL" in auto_name or "ELLIPSE" in auto_name:
            typ = "ellipse"
            radius = 0
        elif "ROUNDED" in auto_name:
            radius = 18
        elif "RECTANGLE" in auto_name:
            radius = 0
    except Exception:
        pass
    return {
        "id": _uid(),
        "type": typ,
        **box,
        "anim": "fade",
        "animDelay": 0,
        "fill": fill,
        "stroke": stroke,
        "strokeW": stroke_w,
        "dashed": _line_is_dashed(shape),
        "radius": radius,
        **_anim_for(shape, anims),
    }


def _line_element(shape, prs, anims):
    box = _box(shape, prs)
    thickness = max(2, _stroke_width(shape, _pt_scale(prs)) or 2)
    if box["w"] >= box["h"]:
        box["h"] = thickness
    else:
        box["w"] = thickness
    return {
        "id": _uid(),
        "type": "line",
        **box,
        "anim": "reveal",
        "animDelay": 0,
        "fill": _stroke_color(shape, "#16140f"),
        "strokeW": thickness,
        "dashed": _line_is_dashed(shape),
        **_anim_for(shape, anims),
    }


def _table_element(shape, prs, anims):
    try:
        tbl = shape.table
        data = []
        for row in tbl.rows:
            data.append([("\n".join(p.text for p in c.text_frame.paragraphs)).strip()
                         for c in row.cells])
        if not data:
            return None
        # Accent from the first header cell's fill, when solid.
        accent = "#1d4e89"
        try:
            first = tbl.rows[0].cells[0]
            if first.fill.type == MSO_FILL.SOLID:
                accent = _color_to_hex(first.fill.fore_color, accent) or accent
        except Exception:
            pass
        return {
            "id": _uid(),
            "type": "table",
            **_box(shape, prs),
            "anim": "rise",
            "animDelay": 0,
            "rows": len(data),
            "cols": max(len(r) for r in data),
            "header": True,
            "accent": accent,
            "theme": "clean",
            "font": '"Archivo",sans-serif',
            "size": 18,
            "tableData": data,
            **_anim_for(shape, anims),
        }
    except Exception:
        return None


_XL_TO_HANNS_KIND = {
    "BAR": "horizontalBar", "BAR_CLUSTERED": "horizontalBar", "BAR_STACKED": "horizontalBar",
    "COLUMN_CLUSTERED": "bar", "COLUMN_STACKED": "stackedBar",
    "PIE": "pie", "PIE_EXPLODED": "pie",
    "DOUGHNUT": "donut", "DOUGHNUT_EXPLODED": "donut",
    "LINE": "line", "LINE_MARKERS": "line", "LINE_STACKED": "line",
    "AREA": "area", "AREA_STACKED": "area",
    "XY_SCATTER": "scatter", "XY_SCATTER_LINES": "scatter", "XY_SCATTER_SMOOTH": "spline",
    "RADAR": "radar", "RADAR_FILLED": "radar", "RADAR_MARKERS": "radar",
}


def _series_color(series, default=None):
    try:
        return _color_to_hex(series.format.fill.fore_color, default)
    except Exception:
        return default


def _chart_element(shape, prs, anims, warnings):
    box = _box(shape, prs)
    title = "Imported chart"
    kind = "bar"
    chart_data = None
    palette = []
    series_names = []
    try:
        chart = shape.chart
        try:
            if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                title = chart.chart_title.text_frame.text or title
        except Exception:
            pass

        xl_name = str(getattr(chart, "chart_type", "") or "")
        xl_name = xl_name.split(" ")[0] if xl_name else ""
        kind = _XL_TO_HANNS_KIND.get(xl_name, "bar")

        plots = list(chart.plots)
        if plots:
            cats = [str(c) for c in plots[0].categories]
            all_series = [s for plot in plots for s in plot.series]
            series_names = [str(s.name or f"Series {i+1}") for i, s in enumerate(all_series)]
            for s in all_series:
                col = _series_color(s)
                if col:
                    palette.append(col)

            if len(all_series) >= 2 and kind in ("bar", "horizontalBar", "stackedBar"):
                # Multi-series bar chart → Hanns grouped/stacked bars.
                if kind != "stackedBar":
                    kind = "groupedBar"
                chart_data = []
                for i, label in enumerate(cats):
                    row = {"label": label, "value": _sv(all_series[0], i)}
                    for j, s in enumerate(all_series[1:], start=2):
                        row[f"value{j}"] = _sv(s, i)
                    chart_data.append(row)
            elif all_series:
                s0 = all_series[0]
                chart_data = [
                    {"label": (cats[i] if i < len(cats) else f"#{i+1}"), "value": _sv(s0, i)}
                    for i in range(max(len(cats), len(list(s0.values) or [])))
                ]
    except Exception as exc:
        warnings.append(f"Chart data could not be fully read ({exc}); imported an editable placeholder.")

    el = {
        "id": _uid(),
        "type": "chart",
        **box,
        "anim": "rise",
        "animDelay": 0,
        "chartKind": kind,
        "chartType": kind,            # legacy key kept for old readers
        "renderEngine": "svg",
        "title": title,
        "accent": (palette[0] if palette else "#e8482b"),
        "showValues": True,
        "showLabels": True,
        "chartData": chart_data or [{"label": "Imported", "value": 1}, {"label": "Edit", "value": 1}],
        **_anim_for(shape, anims),
    }
    if palette:
        el["palette"] = palette + ["#e8482b", "#22c55e", "#38bdf8", "#f59e0b"][: max(0, 4 - len(palette))]
    if len(series_names) > 1:
        el["seriesNames"] = series_names
        el["showLegend"] = True
    return el


def _sv(series, i):
    try:
        vals = list(series.values)
        v = vals[i] if i < len(vals) else 0
        return round(float(v), 3) if v is not None else 0
    except Exception:
        return 0


def _media_element(request, deck, shape, prs, anims):
    """Movies: keep the poster frame as an image so the slide still reads."""
    try:
        el = _image_element(request, deck, shape, prs, anims)
        if el:
            el["alt"] = "Imported video poster frame"
            return el
    except Exception:
        pass
    box = _box(shape, prs)
    return {
        "id": _uid(), "type": "rect", **box,
        "anim": "fade", "animDelay": 0,
        "fill": "#0f172a", "stroke": "none", "strokeW": 0, "radius": 14,
        **_anim_for(shape, anims),
    }


# ─────────────────────────── shape walk ───────────────────────────

def _elements_from_shape(request, deck, shape, prs, warnings, anims):
    elements = []
    stype = getattr(shape, "shape_type", None)

    try:
        if stype == MSO_SHAPE_TYPE.GROUP:
            group_anim = _anim_for(shape, anims)
            for child in getattr(shape, "shapes", []):
                kids = _elements_from_shape(request, deck, child, prs, warnings, anims)
                if group_anim:
                    for k in kids:
                        k.setdefault("anim", group_anim.get("anim", "fade"))
                        if "anim" in group_anim and k.get("anim") == "fade":
                            k.update(group_anim)
                elements.extend(kids)
            return elements
    except Exception:
        pass

    try:
        if getattr(shape, "has_table", False):
            el = _table_element(shape, prs, anims)
            if el:
                elements.append(el)
                return elements
    except Exception:
        pass

    try:
        if getattr(shape, "has_chart", False):
            el = _chart_element(shape, prs, anims, warnings)
            if el:
                elements.append(el)
            return elements
    except Exception:
        pass

    # ANY shape that carries an embedded image is a picture as far as Hanns
    # is concerned. Checking MSO_SHAPE_TYPE.PICTURE alone missed:
    #   • pictures inside a layout's picture/content placeholder
    #     (shape_type == PLACEHOLDER) — the normal way template decks
    #     hold images, and they were silently dropped;
    #   • LINKED_PICTURE shapes;
    #   • picture-filled auto shapes.
    try:
        if _shape_image(shape) is not None:
            el = _image_element(request, deck, shape, prs, anims)
            if el:
                elements.append(el)
                return elements
    except Exception:
        pass

    try:
        if stype == MSO_SHAPE_TYPE.MEDIA:
            el = _media_element(request, deck, shape, prs, anims)
            if el:
                elements.append(el)
            warnings.append("One video was imported as its poster image (Hanns video elements can re-link the file).")
            return elements
    except Exception:
        pass

    try:
        if stype in (MSO_SHAPE_TYPE.LINE,):
            elements.append(_line_element(shape, prs, anims))
            return elements
    except Exception:
        pass

    # Auto shapes: keep the shape AND lift any text into a separate editable
    # text element layered on top (the Hanns shape renderer has no caption).
    is_auto = False
    try:
        is_auto = stype in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM)
    except Exception:
        pass

    has_text = False
    try:
        has_text = bool(getattr(shape, "has_text_frame", False) and _shape_text(shape))
    except Exception:
        pass

    # Divider lines / connectors stored as AUTO_SHAPEs (prst="line",
    # straightConnector1, …). These must become real line elements — the
    # generic auto-shape path would render them as filled 4 px rects.
    if not has_text and _is_line_geometry(shape):
        try:
            elements.append(_line_element(shape, prs, anims))
            return elements
        except Exception:
            pass

    if is_auto:
        # An auto shape whose whole text is one hyperlink is a link button
        # (this is exactly how the Hanns exporter writes link elements).
        if has_text:
            try:
                url = _whole_shape_hyperlink(shape)
                if not url:
                    # Exported Hanns link cards carry the hyperlink on the
                    # label run only; accept a first-paragraph hyperlink too.
                    try:
                        first_p = shape.text_frame.paragraphs[0]
                        for r0 in first_p.runs:
                            addr = getattr(getattr(r0, "hyperlink", None), "address", None)
                            if addr:
                                url = addr
                                break
                    except Exception:
                        url = None
                if url:
                    box = _box(shape, prs)
                    text = _shape_text(shape, with_bullets=False)
                    _p, r = _dominant_run(shape)
                    elements.append({
                        "id": _uid(), "type": "link", **box,
                        "anim": "rise", "animDelay": 0,
                        "url": url,
                        "label": text.split("\n")[0][:80],
                        "description": (text.split("\n")[1][:120] if "\n" in text else url),
                        "linkStyle": "button",
                        "accent": _fill_or(shape, "#2563eb"),
                        "textColor": _font_color(r, "#ffffff"),
                        "bg": _fill_or(shape, "#2563eb"),
                        "radius": 18,
                        **_anim_for(shape, anims),
                    })
                    return elements
            except Exception:
                pass
        try:
            shp_el = _shape_element(shape, prs, anims)
            # Only emit the shape when it is actually visible. Text boxes in
            # many decks are AUTO_SHAPEs with no fill and no outline — the
            # old importer painted every one of them as a beige rectangle
            # behind its text.
            has_visible_fill = shp_el.get("fill") not in (None, "", "none")
            has_visible_stroke = (
                shp_el.get("stroke") not in (None, "", "none")
                and float(shp_el.get("strokeW") or 0) > 0
            )
            if has_visible_fill or has_visible_stroke:
                elements.append(shp_el)
        except Exception:
            pass
        if has_text:
            try:
                txt_el = _text_element(shape, prs, anims)
                if txt_el and txt_el.get("type") == "text":
                    txt_el["fill"] = "none"           # the shape below carries the fill
                    txt_el["align"] = txt_el.get("align") or "center"
                    elements.append(txt_el)
                elif txt_el:
                    elements.append(txt_el)
            except Exception:
                pass
        return elements

    if has_text:
        try:
            el = _text_element(shape, prs, anims)
            if el:
                elements.append(el)
            return elements
        except Exception:
            pass

    return elements


# ─────────────────────────── notes / assembly ───────────────────────────

def _slide_notes(slide, index: int, warnings: list[str]) -> str:
    note = ""
    try:
        if slide.has_notes_slide:
            note = (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        note = ""
    if not note:
        note = f"Imported from PowerPoint slide {index + 1}."
    if warnings:
        note += "\n\n(Some complex PowerPoint features may need manual adjustment in Hanns.)"
    return note


def import_powerpoint_into_deck(*, request, deck, uploaded_file, replace: bool = True) -> dict:
    prs, original_name, tmpdir = _load_presentation(uploaded_file)
    warnings: list[str] = []

    try:
        stem = Path(original_name).stem.strip()
        if stem:
            deck.title = stem[:140]
            deck.save(update_fields=["title", "updated_at"])

        if replace:
            deck.slides.all().delete()

        slide_rows = []
        for idx, slide in enumerate(prs.slides):
            slide_warnings: list[str] = []
            anims = _entrance_animations(slide)
            elements = []
            for shape in slide.shapes:
                try:
                    elements.extend(
                        _elements_from_shape(request, deck, shape, prs, slide_warnings, anims)
                    )
                except Exception as exc:
                    slide_warnings.append(f"Skipped one unsupported object on slide {idx + 1}: {exc}")

            if not elements:
                elements.append({
                    "id": _uid(), "type": "text",
                    "x": 90, "y": 210, "w": 780, "h": 90, "rot": 0,
                    "anim": "fade", "animDelay": 0,
                    "text": f"Imported slide {idx + 1}",
                    "font": '"Fraunces",serif', "size": 44, "weight": 700,
                    "italic": False, "color": "#16140f", "align": "center",
                    "lh": 1.15, "ls": 0, "fill": "none",
                })
                slide_warnings.append(f"Slide {idx + 1} had no editable objects Hanns could import.")

            bg_css, bg_size = _slide_background(request, deck, slide)

            warnings.extend(slide_warnings)
            slide_rows.append(Slide(deck=deck, position=idx, data={
                "bg": bg_css,
                "bgSize": bg_size,
                "bgFx": "none",
                "transition": "fade",
                "notes": _slide_notes(slide, idx, slide_warnings),
                "els": elements,
            }))

        if slide_rows:
            Slide.objects.bulk_create(slide_rows)
        else:
            Slide.objects.create(deck=deck, position=0, data={
                "bg": "#f6f1e7", "bgSize": None, "bgFx": "none",
                "transition": "fade",
                "notes": "PowerPoint import created a blank starter slide.",
                "els": [],
            })
            warnings.append("The PowerPoint file did not contain any slides.")

        return {"warnings": warnings, "slide_count": len(slide_rows)}
    finally:
        tmpdir.cleanup()