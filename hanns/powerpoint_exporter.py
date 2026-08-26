"""Hanns JSON → PowerPoint exporter (v2 — high-fidelity).

The inverse of ``powerpoint_importer.py``. It walks a Deck's slides and their
element JSON ({bg, bgSize, transition, notes, els:[…]}) and writes a .pptx
using python-pptx, keeping the exported file as close as possible to what the
live Hanns stage shows.

What v2 exports beyond the original exporter:

  • Slide backgrounds — solid colours AND real PPTX gradient fills built from
    the deck's CSS gradients (stops + angle), plus full-bleed background
    pictures for ``url(…)`` backgrounds.
  • Text             — top anchor (matches the renderer), line height,
    letter spacing, per-line paragraphs, fonts, bold/italic/colour/align,
    optional box fill.
  • Shapes           — proportional rounded-corner radius, dashed strokes,
    rotation.
  • Images           — ``fit:"contain"`` letterboxes inside the box and
    ``fit:"cover"`` crops, so pictures land exactly like the stage.
  • Tables           — Hanns ``table`` elements become native PPTX tables with
    a styled header row.
  • Charts           — reads the live ``chartKind`` key (``chartType`` kept as
    a fallback), maps bar / horizontal bar / grouped / stacked / line /
    spline / area / pie / donut / scatter, exports EVERY series, applies the
    element palette per-series (and per-point for pie/donut).
  • Links & videos   — become click-to-open shapes with real hyperlinks.
  • Rich elements    — maps / objects / creative shapes export as a styled,
    labelled card rather than disappearing.
  • ANIMATIONS       — every element's Hanns entrance (fade, rise, drop,
    left, right, zoom, pop, blur, reveal, bounce, elastic, flip, spin, skew,
    typewriter, float, …) is written into the slide's native ``<p:timing>``
    tree as the matching PowerPoint entrance effect, preserving each
    element's delay and duration. One click plays the whole build, exactly
    like the Hanns stage.

Design space is 960×540 (16:9), matching the importer's DESIGN_W/H. We render
the PPTX at a real 10in × 5.625in so 1 design px = 1/96 in.
"""

from __future__ import annotations

import io
import re
import os
from urllib.parse import urlparse, unquote

from django.core.files.storage import default_storage

try:  # python-pptx is optional until export is used.
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn, nsdecls
except Exception:  # pragma: no cover - handled at runtime
    Presentation = None

try:
    from PIL import Image as PILImage
except Exception:  # pragma: no cover
    PILImage = None


DESIGN_W = 960
DESIGN_H = 540

# 960×540 design px → 10in × 5.625in slide. 96 px per inch.
PX_PER_IN = 96.0
EMU_PER_IN = 914400


# ───────────────────────── unit / colour helpers ─────────────────────────

def _px_to_emu(px) -> int:
    try:
        return int(round((float(px) / PX_PER_IN) * EMU_PER_IN))
    except Exception:
        return 0


def _clean_hex(value):
    """Return a 6-char uppercase hex (no #) or None for non-solid colours."""
    if value is None:
        return None
    s = str(value).strip()
    if not s or s.lower() in ("none", "transparent"):
        return None
    if "gradient" in s.lower():
        m = re.search(r"#([0-9a-fA-F]{6}|[0-9a-fA-F]{3})", s)
        if not m:
            return None
        s = "#" + m.group(1)
    m = re.match(r"^#?([0-9a-fA-F]{6})$", s)
    if m:
        return m.group(1).upper()
    m = re.match(r"^#?([0-9a-fA-F]{3})$", s)
    if m:
        r, g, b = m.group(1)
        return (r + r + g + g + b + b).upper()
    m = re.match(r"^rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)", s)
    if m:
        return "{:02X}{:02X}{:02X}".format(
            min(255, int(m.group(1))), min(255, int(m.group(2))), min(255, int(m.group(3)))
        )
    return None


def _rgb(value, default="000000"):
    hexv = _clean_hex(value) or default
    try:
        return RGBColor.from_string(hexv)
    except Exception:
        return RGBColor.from_string(default)


def _num(value, default=0.0):
    try:
        return float(value)
    except Exception:
        return default


# ───────────────────────── CSS gradient parsing ─────────────────────────

def _split_top_level(css: str) -> list[str]:
    """Split layered CSS backgrounds on top-level commas."""
    parts, depth, cur = [], 0, ""
    for ch in css:
        if ch == "(":
            depth += 1
        elif ch == ")":
            depth -= 1
        if ch == "," and depth == 0:
            parts.append(cur.strip())
            cur = ""
        else:
            cur += ch
    if cur.strip():
        parts.append(cur.strip())
    return parts


_COLOR_STOP_RE = re.compile(
    r"(#[0-9a-fA-F]{3,6}|rgba?\([^)]*\))\s*([0-9.]+%)?"
)


def _parse_css_gradient(css: str):
    """Return (angle_deg, [(pos_0_100, '#RRGGBB'), …]) or None.

    Layered backgrounds pick the LAST layer (the base coat) — that is what
    dominates visually and what a projector fallback should show.
    """
    if not css or "gradient" not in css:
        return None
    layers = [p for p in _split_top_level(css) if "gradient" in p]
    if not layers:
        return None
    layer = layers[-1]

    angle = 135.0
    m = re.search(r"linear-gradient\(\s*([0-9.]+)deg", layer)
    if m:
        angle = float(m.group(1))

    inner = layer[layer.find("(") + 1: layer.rfind(")")]
    stops = []
    for cm in _COLOR_STOP_RE.finditer(inner):
        hexv = _clean_hex(cm.group(1))
        if not hexv:
            continue
        pos = None
        if cm.group(2):
            try:
                pos = float(cm.group(2).rstrip("%"))
            except ValueError:
                pos = None
        stops.append([pos, "#" + hexv])
    if len(stops) < 2:
        return None
    # Fill in missing positions evenly.
    if stops[0][0] is None:
        stops[0][0] = 0.0
    if stops[-1][0] is None:
        stops[-1][0] = 100.0
    known = [i for i, s in enumerate(stops) if s[0] is not None]
    for a, b in zip(known, known[1:]):
        span = stops[b][0] - stops[a][0]
        gap = b - a
        for k in range(a + 1, b):
            stops[k][0] = stops[a][0] + span * (k - a) / gap
    stops = [(max(0.0, min(100.0, p)), c) for p, c in stops]
    stops.sort(key=lambda s: s[0])
    return angle, stops


def _bg_url(css: str) -> str | None:
    m = re.search(r"""url\(\s*['"]?([^'")]+)['"]?\s*\)""", css or "")
    return m.group(1) if m else None


# ───────────────────────── image fetch ─────────────────────────

def _load_image_bytes(src):
    """Return (BytesIO, ext) for an element image src, or (None, None).

    Handles a stored media URL, a bare storage-relative path, or a data: URI.
    Network URLs are intentionally NOT fetched (keeps export hermetic and
    avoids SSRF); unresolved images become a placeholder box instead.
    """
    if not src:
        return None, None
    s = str(src)

    if s.startswith("data:"):
        try:
            import base64
            header, b64 = s.split(",", 1)
            ext = "png"
            m = re.search(r"data:image/(\w+)", header)
            if m:
                ext = m.group(1).lower()
                if ext == "jpeg":
                    ext = "jpg"
            return io.BytesIO(base64.b64decode(b64)), ext
        except Exception:
            return None, None

    path = s
    if s.startswith("http://") or s.startswith("https://"):
        path = unquote(urlparse(s).path)
    for prefix in ("/media/", "media/"):
        idx = path.find(prefix)
        if idx != -1:
            path = path[idx + len(prefix):]
            break
    path = path.lstrip("/")

    ext = (os.path.splitext(path)[1] or ".png").lstrip(".").lower()
    if ext == "jpeg":
        ext = "jpg"

    try:
        if default_storage.exists(path):
            with default_storage.open(path, "rb") as fh:
                return io.BytesIO(fh.read()), ext
    except Exception:
        pass
    return None, None


def _image_size(stream) -> tuple[int, int] | None:
    if PILImage is None or stream is None:
        return None
    try:
        pos = stream.tell()
        img = PILImage.open(stream)
        size = img.size
        stream.seek(pos)
        return size
    except Exception:
        try:
            stream.seek(0)
        except Exception:
            pass
        return None


# ───────────────────────── animation registry ─────────────────────────

class _AnimFx:
    """Collects (shape_id, anim, delay_s, dur_s) per slide, then writes the
    native <p:timing> tree."""

    # Hanns anim → (presetID, presetSubtype, behaviour builder key, default ms)
    MAP = {
        "fade":       (10, 0, "fade", 620),
        "blur":       (10, 0, "fade", 620),
        "rise":       (2, 4, "fly_bottom", 620),
        "float":      (2, 4, "fly_bottom", 760),
        "bounce":     (2, 4, "fly_bottom", 950),
        "drop":       (2, 1, "fly_top", 620),
        "left":       (2, 8, "fly_left", 620),
        "skew":       (2, 8, "fly_left", 620),
        "right":      (2, 2, "fly_right", 620),
        "zoom":       (23, 16, "zoom", 620),
        "pop":        (23, 16, "zoom", 720),
        "elastic":    (23, 16, "zoom", 980),
        "blurzoom":   (23, 16, "zoom", 820),
        "reveal":     (22, 8, "wipe_left", 620),
        "typewriter": (22, 8, "wipe_left", 900),
        "revealUp":   (22, 4, "wipe_up", 620),
        "spin":       (21, 1, "spin", 760),
        "flipx":      (19, 0, "flip_x", 680),
        "flipy":      (19, 0, "flip_y", 680),
    }

    def __init__(self):
        self.effects: list[tuple[int, str, float, float]] = []
        self._id = 2  # 1=tmRoot, 2=mainSeq

    def register(self, shape, el):
        anim = str(el.get("anim") or "none")
        if anim == "none" or anim not in self.MAP:
            return
        spid = getattr(shape, "shape_id", None)
        if spid is None:
            try:
                spid = int(shape._element.nvSpPr.cNvPr.get("id"))
            except Exception:
                return
        delay = max(0.0, _num(el.get("animDelay"), 0.0))
        dur = _num(el.get("animDur"), 0.0)
        if dur <= 0:
            dur = self.MAP[anim][3] / 1000.0
        self.effects.append((int(spid), anim, delay, dur))

    # ── XML builders ──
    def _nid(self) -> int:
        self._id += 1
        return self._id

    def _bhv_fade(self, spid, dur):
        return (
            f'<p:animEffect transition="in" filter="fade">'
            f'<p:cBhvr><p:cTn id="{self._nid()}" dur="{dur}"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
        )

    def _bhv_wipe(self, spid, dur, direction):
        return (
            f'<p:animEffect transition="in" filter="wipe({direction})">'
            f'<p:cBhvr><p:cTn id="{self._nid()}" dur="{dur}"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr></p:animEffect>'
        )

    def _bhv_anim_var(self, spid, dur, attr, frm):
        return (
            f'<p:anim calcmode="lin" valueType="num">'
            f'<p:cBhvr additive="base"><p:cTn id="{self._nid()}" dur="{dur}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>{attr}</p:attrName></p:attrNameLst></p:cBhvr>'
            f'<p:tavLst>'
            f'<p:tav tm="0"><p:val><p:strVal val="{frm}"/></p:val></p:tav>'
            f'<p:tav tm="100000"><p:val><p:strVal val="#{attr}"/></p:val></p:tav>'
            f'</p:tavLst></p:anim>'
        )

    def _bhv_fly(self, spid, dur, direction):
        if direction == "bottom":
            return (self._bhv_anim_var(spid, dur, "ppt_x", "#ppt_x")
                    + self._bhv_anim_var(spid, dur, "ppt_y", "1+#ppt_h/2"))
        if direction == "top":
            return (self._bhv_anim_var(spid, dur, "ppt_x", "#ppt_x")
                    + self._bhv_anim_var(spid, dur, "ppt_y", "0-#ppt_h/2"))
        if direction == "left":
            return (self._bhv_anim_var(spid, dur, "ppt_x", "0-#ppt_w/2")
                    + self._bhv_anim_var(spid, dur, "ppt_y", "#ppt_y"))
        return (self._bhv_anim_var(spid, dur, "ppt_x", "1+#ppt_w/2")
                + self._bhv_anim_var(spid, dur, "ppt_y", "#ppt_y"))

    def _bhv_scale(self, spid, dur, fx, fy):
        return (
            f'<p:animScale><p:cBhvr><p:cTn id="{self._nid()}" dur="{dur}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
            f'<p:from x="{fx}" y="{fy}"/><p:to x="100000" y="100000"/></p:animScale>'
        )

    def _bhv_spin(self, spid, dur):
        return (
            f'<p:animRot by="21600000"><p:cBhvr><p:cTn id="{self._nid()}" dur="{dur}" fill="hold"/>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            f'<p:attrNameLst><p:attrName>r</p:attrName></p:attrNameLst></p:cBhvr></p:animRot>'
        )

    def _behaviours(self, key, spid, dur):
        if key == "fade":
            return self._bhv_fade(spid, dur)
        if key == "fly_bottom":
            return self._bhv_fade(spid, dur) + self._bhv_fly(spid, dur, "bottom")
        if key == "fly_top":
            return self._bhv_fade(spid, dur) + self._bhv_fly(spid, dur, "top")
        if key == "fly_left":
            return self._bhv_fade(spid, dur) + self._bhv_fly(spid, dur, "left")
        if key == "fly_right":
            return self._bhv_fade(spid, dur) + self._bhv_fly(spid, dur, "right")
        if key == "zoom":
            return self._bhv_fade(spid, dur) + self._bhv_scale(spid, dur, 10000, 10000)
        if key == "wipe_left":
            return self._bhv_wipe(spid, dur, "left")
        if key == "wipe_up":
            return self._bhv_wipe(spid, dur, "up")
        if key == "spin":
            return self._bhv_fade(spid, dur) + self._bhv_spin(spid, dur)
        if key == "flip_x":
            return self._bhv_fade(spid, dur) + self._bhv_scale(spid, dur, 1000, 100000)
        if key == "flip_y":
            return self._bhv_fade(spid, dur) + self._bhv_scale(spid, dur, 100000, 1000)
        return self._bhv_fade(spid, dur)

    def apply(self, slide):
        """Append the <p:timing> tree to ``slide`` (a python-pptx Slide)."""
        if not self.effects:
            return
        effect_xml = []
        for i, (spid, anim, delay, dur_s) in enumerate(self.effects):
            preset_id, subtype, key, _default = self.MAP[anim]
            dur = max(1, int(round(dur_s * 1000)))
            delay_ms = max(0, int(round(delay * 1000)))
            node_type = "clickEffect" if i == 0 else "withEffect"
            ctn_id = self._nid()
            set_id = self._nid()
            effect_xml.append(
                f'<p:par><p:cTn id="{ctn_id}" presetID="{preset_id}" presetClass="entr" '
                f'presetSubtype="{subtype}" fill="hold" grpId="0" nodeType="{node_type}">'
                f'<p:stCondLst><p:cond delay="{delay_ms}"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:set><p:cBhvr><p:cTn id="{set_id}" dur="1" fill="hold">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
                f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
                f'{self._behaviours(key, spid, dur)}'
                f'</p:childTnLst></p:cTn></p:par>'
            )
        bld = "".join(
            f'<p:bldP spid="{spid}" grpId="0"/>'
            for spid in {e[0] for e in self.effects}
        )
        xml = (
            f'<p:timing {nsdecls("p", "a")}>'
            f'<p:tnLst><p:par>'
            f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
            f'<p:seq concurrent="1" nextAc="seek">'
            f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>'
            f'<p:par><p:cTn id="{self._nid()}" fill="hold">'
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>'
            f'<p:par><p:cTn id="{self._nid()}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>'
            f'{"".join(effect_xml)}'
            f'</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn></p:par>'
            f'</p:childTnLst></p:cTn>'
            f'<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
            f'<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
            f'</p:seq>'
            f'</p:childTnLst></p:cTn>'
            f'</p:par></p:tnLst>'
            f'<p:bldLst>{bld}</p:bldLst>'
            f'</p:timing>'
        )
        try:
            slide._element.append(parse_xml(xml))
        except Exception:
            pass  # timing is a bonus — never abort the export over it


# ───────────────────────── element renderers ─────────────────────────

def _css_font_to_name(css):
    """'"Inter",sans-serif' → 'Inter'."""
    if not css:
        return None
    first = str(css).split(",")[0].strip().strip('"').strip("'")
    return first or None


def _set_letter_spacing(run, px):
    """Letter spacing: Hanns stores px, PPTX wants 1/100 pt on <a:rPr spc>."""
    try:
        pts = float(px) / 1.333
        run.font._rPr.set("spc", str(int(round(pts * 100))))
    except Exception:
        pass


def _add_text(slide, el):
    box = slide.shapes.add_textbox(
        _px_to_emu(el.get("x", 0)), _px_to_emu(el.get("y", 0)),
        _px_to_emu(max(8, el.get("w", 100))), _px_to_emu(max(8, el.get("h", 40))),
    )
    tf = box.text_frame
    tf.word_wrap = True
    try:
        # The Hanns renderer top-aligns text inside its box.
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Emu(int(0.02 * EMU_PER_IN))
        tf.margin_top = tf.margin_bottom = Emu(int(0.01 * EMU_PER_IN))
    except Exception:
        pass

    fill_hex = _clean_hex(el.get("fill"))
    if fill_hex:
        try:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
            box.line.fill.background()
        except Exception:
            pass

    text = str(el.get("text", "") or "")
    lines = text.split("\n")
    align = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
        str(el.get("align", "left")).lower(), PP_ALIGN.LEFT
    )
    size_pt = max(6, int(round(_num(el.get("size", 24), 24) / 1.333)))
    color = _rgb(el.get("color", "#16140f"), "16140F")
    bold = _num(el.get("weight", 500), 500) >= 600
    italic = bool(el.get("italic", False))
    font_name = _css_font_to_name(el.get("font"))
    lh = _num(el.get("lh", 1.15), 1.15)
    ls = _num(el.get("ls", 0), 0)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try:
            p.line_spacing = max(0.5, min(3.0, lh))
        except Exception:
            pass
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(size_pt)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if font_name:
            f.name = font_name
        if ls:
            _set_letter_spacing(run, ls)
    rot = _num(el.get("rot", 0), 0)
    if rot:
        box.rotation = rot
    return box


def _add_image(slide, el):
    stream, _ext = _load_image_bytes(el.get("src"))
    x, y = _px_to_emu(el.get("x", 0)), _px_to_emu(el.get("y", 0))
    w, h = _px_to_emu(max(8, el.get("w", 100))), _px_to_emu(max(8, el.get("h", 100)))
    fit = str(el.get("fit", "cover")).lower()

    if stream is not None:
        try:
            size = _image_size(stream)
            if size and fit == "contain":
                iw, ih = size
                scale = min(w / iw, h / ih)
                dw, dh = int(iw * scale), int(ih * scale)
                pic = slide.shapes.add_picture(
                    stream, x + (w - dw) // 2, y + (h - dh) // 2, width=dw, height=dh
                )
            else:
                pic = slide.shapes.add_picture(stream, x, y, width=w, height=h)
                if size and fit == "cover":
                    iw, ih = size
                    box_ar, img_ar = w / h, iw / ih
                    if img_ar > box_ar:      # image too wide → crop sides
                        crop = (1 - box_ar / img_ar) / 2
                        pic.crop_left = crop
                        pic.crop_right = crop
                    elif img_ar < box_ar:    # image too tall → crop top/bottom
                        crop = (1 - img_ar / box_ar) / 2
                        pic.crop_top = crop
                        pic.crop_bottom = crop
            rot = _num(el.get("rot", 0), 0)
            if rot:
                pic.rotation = rot
            return pic
        except Exception:
            pass

    # Placeholder when the image couldn't be resolved.
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string("E7E8D1")
    shp.line.color.rgb = RGBColor.from_string("B85042")
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Image"
    r.font.size = Pt(12)
    r.font.color.rgb = RGBColor.from_string("6B6354")
    return shp


def _set_rounded_radius(shp, radius_px, box_h_px):
    """Match the CSS px radius: adjustment ≈ radius / (shorter side / 2)."""
    try:
        frac = max(0.0, min(0.5, float(radius_px) / max(1.0, float(box_h_px))))
        shp.adjustments[0] = frac
    except Exception:
        pass


def _set_dashed(shp):
    try:
        ln = shp.line._get_or_add_ln()
        dash = ln.find(qn("a:prstDash"))
        if dash is None:
            dash = parse_xml(f'<a:prstDash {nsdecls("a")} val="dash"/>')
            ln.append(dash)
        else:
            dash.set("val", "dash")
    except Exception:
        pass


def _add_shape(slide, el, kind):
    x, y = _px_to_emu(el.get("x", 0)), _px_to_emu(el.get("y", 0))
    w, h = _px_to_emu(max(4, el.get("w", 40))), _px_to_emu(max(4, el.get("h", 40)))
    radius = _num(el.get("radius", 0), 0)
    if kind == "ellipse":
        auto = MSO_SHAPE.OVAL
    elif radius > 0:
        auto = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        auto = MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(auto, x, y, w, h)
    if auto == MSO_SHAPE.ROUNDED_RECTANGLE:
        _set_rounded_radius(shp, radius, min(_num(el.get("w", 40)), _num(el.get("h", 40))))

    fill_hex = _clean_hex(el.get("fill"))
    if fill_hex:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    else:
        shp.fill.background()

    stroke_hex = _clean_hex(el.get("stroke"))
    stroke_w = _num(el.get("strokeW", 0), 0)
    if stroke_hex and stroke_w > 0:
        shp.line.color.rgb = RGBColor.from_string(stroke_hex)
        shp.line.width = Pt(max(0.5, stroke_w / 1.333))
        if el.get("dashed"):
            _set_dashed(shp)
    else:
        shp.line.fill.background()

    rot = _num(el.get("rot", 0), 0)
    if rot:
        shp.rotation = rot

    txt = el.get("text")
    if txt:
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(txt)
        r.font.size = Pt(max(8, int(round(_num(el.get("size", 18), 18) / 1.333))))
        r.font.color.rgb = _rgb(el.get("color", "#16140f"), "16140F")
    return shp


def _add_line(slide, el):
    """Lines are stored as a thin box; draw a real PPTX connector."""
    x = _px_to_emu(el.get("x", 0))
    y = _px_to_emu(el.get("y", 0))
    w = _px_to_emu(max(2, el.get("w", 100)))
    h = _px_to_emu(max(1, el.get("h", 2)))
    try:
        from pptx.enum.shapes import MSO_CONNECTOR
        vertical = _num(el.get("h", 2)) > _num(el.get("w", 100))
        if vertical:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + w // 2, y, x + w // 2, y + h)
            width_pt = max(0.75, _num(el.get("w", 2)) / 1.333)
        else:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y + h // 2, x + w, y + h // 2)
            width_pt = max(0.75, _num(el.get("h", 2)) / 1.333)
        conn.line.color.rgb = _rgb(el.get("fill", "#16140f"), "16140F")
        conn.line.width = Pt(min(20.0, width_pt))
        if el.get("dashed"):
            _set_dashed(conn)
        return conn
    except Exception:
        pass
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, max(_px_to_emu(2), h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(el.get("fill", "#16140f"), "16140F")
    shp.line.fill.background()
    return shp


# ───────────────────────── tables ─────────────────────────

def _add_table(slide, el):
    data = el.get("tableData") or []
    data = [r for r in data if isinstance(r, (list, tuple))]
    if not data:
        return None
    rows = len(data)
    cols = max(len(r) for r in data)
    x, y = _px_to_emu(el.get("x", 0)), _px_to_emu(el.get("y", 0))
    w, h = _px_to_emu(max(60, el.get("w", 400))), _px_to_emu(max(40, el.get("h", 200)))

    gframe = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = gframe.table
    accent = _clean_hex(el.get("accent")) or "1D4E89"
    size_pt = max(8, int(round(_num(el.get("size", 18), 18) / 1.333)))
    font_name = _css_font_to_name(el.get("font")) or "Archivo"
    header = el.get("header", True)

    for r in range(rows):
        row_data = data[r]
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(row_data[c]) if c < len(row_data) else ""
            para = cell.text_frame.paragraphs[0]
            for run in (para.runs or [para.add_run()]):
                run.font.size = Pt(size_pt)
                run.font.name = font_name
                if r == 0 and header:
                    run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string("FFFFFF")
            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(accent)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string("F4F6F8")
    return gframe


# ───────────────────────── charts ─────────────────────────

_KIND_TO_XL = {
    "bar": "COLUMN_CLUSTERED",
    "column": "COLUMN_CLUSTERED",
    "horizontalbar": "BAR_CLUSTERED",
    "groupedbar": "COLUMN_CLUSTERED",
    "stackedbar": "COLUMN_STACKED",
    "pie": "PIE",
    "donut": "DOUGHNUT",
    "doughnut": "DOUGHNUT",
    "line": "LINE_MARKERS",
    "spline": "LINE_MARKERS",
    "area": "AREA",
    "radar": "RADAR",
    "scatter": "XY_SCATTER",
    "bubble": "XY_SCATTER",
    # Kinds with no PPTX equivalent degrade to a clustered column chart so
    # the data itself always survives the export.
    "gauge": "DOUGHNUT",
    "progress": "BAR_CLUSTERED",
    "funnel": "BAR_CLUSTERED",
    "waterfall": "COLUMN_CLUSTERED",
    "heatmap": "COLUMN_CLUSTERED",
    "treemap": "COLUMN_CLUSTERED",
    "kpi": "COLUMN_CLUSTERED",
}


def _chart_series(el):
    """Return (categories, [(name, values), …]) from the element's data."""
    data = el.get("chartData") or []
    cats, base_vals = [], []
    extra_cols: dict[int, list] = {}
    for row in data:
        if not isinstance(row, dict):
            continue
        cats.append(str(row.get("label", "")))
        base_vals.append(_num(row.get("value", 0), 0))
        for key, val in row.items():
            m = re.match(r"^value(\d+)$", str(key))
            if m:
                extra_cols.setdefault(int(m.group(1)), []).append(_num(val, 0))
    if not cats:
        cats, base_vals = ["A", "B"], [1, 1]

    names = el.get("seriesNames") or []
    series = [(str(names[0]) if names else str(el.get("title") or "Series 1"), base_vals)]
    for idx in sorted(extra_cols):
        vals = extra_cols[idx]
        vals += [0] * (len(cats) - len(vals))
        label = str(names[idx - 1]) if len(names) >= idx else f"Series {idx}"
        series.append((label, vals[: len(cats)]))
    return cats, series


def _add_chart(slide, el):
    kind = str(el.get("chartKind") or el.get("chartType") or "bar").lower()
    xl_name = _KIND_TO_XL.get(kind, "COLUMN_CLUSTERED")
    xl_type = getattr(XL_CHART_TYPE, xl_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

    cats, series = _chart_series(el)
    chart_data = CategoryChartData()
    chart_data.categories = cats
    for name, vals in series:
        chart_data.add_series(name, vals)

    x, y = _px_to_emu(el.get("x", 0)), _px_to_emu(el.get("y", 0))
    w, h = _px_to_emu(max(80, el.get("w", 320))), _px_to_emu(max(80, el.get("h", 220)))

    try:
        gframe = slide.shapes.add_chart(xl_type, x, y, w, h, chart_data)
        chart = gframe.chart

        multi = len(series) > 1
        chart.has_legend = bool(el.get("showLegend", multi) or kind in ("pie", "donut", "doughnut"))
        if chart.has_legend:
            try:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
            except Exception:
                pass

        if el.get("title"):
            try:
                chart.has_title = True
                chart.chart_title.text_frame.text = str(el.get("title"))
            except Exception:
                pass

        if el.get("showValues", True):
            try:
                plot = chart.plots[0]
                plot.has_data_labels = True
                plot.data_labels.number_format = "0.##"
                plot.data_labels.number_format_is_linked = False
            except Exception:
                pass

        palette = [
            _clean_hex(c) for c in (el.get("palette") or [])
            if _clean_hex(c)
        ] or [c for c in [_clean_hex(el.get("accent"))] if c]
        if palette:
            try:
                if kind in ("pie", "donut", "doughnut") and chart.series:
                    # Colour each slice from the palette.
                    s0 = chart.series[0]
                    for i, point in enumerate(s0.points):
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = RGBColor.from_string(
                            palette[i % len(palette)]
                        )
                else:
                    for i, s in enumerate(chart.series):
                        s.format.fill.solid()
                        s.format.fill.fore_color.rgb = RGBColor.from_string(
                            palette[i % len(palette)]
                        )
            except Exception:
                pass
        return gframe
    except Exception:
        return _add_shape(slide, {**el, "type": "rect", "fill": "#ECE4D4",
                                  "text": (el.get("title") or "Chart")}, "rect")


# ───────────────────────── links / video / rich placeholders ─────────────

def _add_hyperlink_card(slide, el, *, label, sub, url, fill, text_color):
    x, y = _px_to_emu(el.get("x", 0)), _px_to_emu(el.get("y", 0))
    w, h = _px_to_emu(max(40, el.get("w", 200))), _px_to_emu(max(24, el.get("h", 60)))
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_rounded_radius(shp, _num(el.get("radius", 18), 18),
                        min(_num(el.get("w", 200)), _num(el.get("h", 60))))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(fill, "2563EB")
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(label or "Open link")
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = _rgb(text_color, "FFFFFF")
    if url:
        try:
            r.hyperlink.address = str(url)
        except Exception:
            pass
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = str(sub)[:80]
        r2.font.size = Pt(10)
        r2.font.color.rgb = _rgb(text_color, "FFFFFF")
    return shp


def _add_link(slide, el):
    return _add_hyperlink_card(
        slide, el,
        label=el.get("label") or "Open link",
        sub=el.get("description") or el.get("url") or "",
        url=el.get("url"),
        fill=el.get("bg") or el.get("accent") or "#2563eb",
        text_color=el.get("textColor") or "#ffffff",
    )


def _add_video(slide, el):
    return _add_hyperlink_card(
        slide, el,
        label="▶ " + (str(el.get("title") or "Play video")),
        sub=el.get("src") or "",
        url=el.get("src"),
        fill="#0f172a",
        text_color="#ffffff",
    )


def _add_rich_placeholder(slide, el, kind_label):
    """Maps / objects / creative shapes: a styled labelled card so nothing
    silently vanishes from the exported deck."""
    label = (el.get("title") or el.get("label") or el.get("statTitle")
             or el.get("nodeTitle") or kind_label)
    shp = _add_shape(slide, {
        **el, "type": "rect",
        "fill": el.get("fill") if _clean_hex(el.get("fill")) else "#F1F5F9",
        "stroke": el.get("accent") or "#94A3B8", "strokeW": 1.5,
        "radius": 18, "text": f"{label}",
        "color": "#0F172A", "size": 20,
    }, "rect")
    return shp


# ───────────────────────── slide background ─────────────────────────

def _set_gradient_bg(slide, angle_deg, stops):
    """Write a real <a:gradFill> background from parsed CSS stops."""
    try:
        gs_xml = "".join(
            f'<a:gs pos="{int(round(p * 1000))}"><a:srgbClr val="{_clean_hex(c)}"/></a:gs>'
            for p, c in stops if _clean_hex(c)
        )
        # CSS angle is clockwise from 12 o'clock; PPTX from 3 o'clock (60000ths).
        ppt_ang = int(round(((angle_deg - 90) % 360) * 60000))
        bg_xml = (
            f'<p:bg {nsdecls("p", "a")}><p:bgPr>'
            f'<a:gradFill rotWithShape="1"><a:gsLst>{gs_xml}</a:gsLst>'
            f'<a:lin ang="{ppt_ang}" scaled="1"/></a:gradFill>'
            f'<a:effectLst/></p:bgPr></p:bg>'
        )
        cSld = slide._element.find(qn("p:cSld"))
        old = cSld.find(qn("p:bg"))
        if old is not None:
            cSld.remove(old)
        cSld.insert(0, parse_xml(bg_xml))
        return True
    except Exception:
        return False


def _set_slide_bg(slide, slide_data):
    bg = str(slide_data.get("bg") or "")

    # 1) Full-bleed background picture.
    url = _bg_url(bg)
    if url:
        stream, _ext = _load_image_bytes(url)
        if stream is not None:
            try:
                slide.shapes.add_picture(stream, 0, 0,
                                         width=_px_to_emu(DESIGN_W),
                                         height=_px_to_emu(DESIGN_H))
                return
            except Exception:
                pass

    # 2) Real gradient fill.
    grad = _parse_css_gradient(bg)
    if grad and _set_gradient_bg(slide, grad[0], grad[1]):
        return

    # 3) Solid colour (or first colour of an unparseable gradient).
    hexv = _clean_hex(bg)
    if not hexv:
        return
    try:
        b = slide.background
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor.from_string(hexv)
    except Exception:
        pass


def _set_notes(slide, slide_data):
    parts = []
    if slide_data.get("notes"):
        parts.append(str(slide_data["notes"]))
    transition = slide_data.get("transition")
    if transition and str(transition).lower() not in ("", "none", "fade"):
        parts.append(f"[Hanns transition: {transition}]")
    if not parts:
        return
    try:
        slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)
    except Exception:
        pass


# ───────────────────────── element dispatch ─────────────────────────

# ─────────────────── studio objects (data & diagram) ───────────────────
#
# Studio objects are drawn in the browser as SVG. PowerPoint has no SVG
# canvas worth targeting, so each one is rebuilt from NATIVE shapes and
# text boxes: a KPI tile becomes a rounded rectangle, a ranked bar becomes
# two rectangles, a process step becomes a real PPTX chevron. The result
# is editable in PowerPoint rather than a flat picture.
#
# The geometric ones (choropleth, Sankey, Venn, 2x2, slope, pyramid) have
# no honest native equivalent, so they fall back to a card that at least
# carries the title and the data as text — nothing silently vanishes and
# the numbers survive the round trip.

STUDIO_KINDS = {
    "choropleth", "gradient_legend", "stat_block", "kpi_grid", "bullet_bars",
    "slope_chart", "waffle", "ring_grid", "process_steps", "timeline_track",
    "rank_bars", "matrix_2x2", "venn", "pyramid_tiers", "sankey_flow",
    "heat_grid", "quote_card",
    # mechanical objects
    "sand_timer", "clock_face", "gears", "charge_meter", "temp_gauge", "speedometer",
}

# Sequential ramps, kept in step with RAMPS in hanns_core.js.
_STUDIO_RAMPS = {
    "ocean":  ["E0F2FE", "7DD3FC", "38BDF8", "0284C7", "075985", "0C3A5B"],
    "teal":   ["CCFBF1", "5EEAD4", "2DD4BF", "0D9488", "115E59", "0B3B38"],
    "forest": ["DCFCE7", "86EFAC", "4ADE80", "16A34A", "15803D", "14532D"],
    "sunset": ["FEF3C7", "FCD34D", "FB923C", "EF4444", "B91C1C", "7F1D1D"],
    "ember":  ["FFEDD5", "FDBA74", "FB923C", "EA580C", "C2410C", "7C2D12"],
    "violet": ["EDE9FE", "C4B5FD", "A78BFA", "7C3AED", "5B21B6", "3B0764"],
    "slate":  ["F1F5F9", "CBD5E1", "94A3B8", "64748B", "334155", "0F172A"],
    "clay":   ["F6F1E7", "E3D5BD", "C9AD82", "A5834F", "6F5533", "3F2F1C"],
}


def _mix_hex(a, b, t):
    try:
        ar, ag, ab = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
        br, bg, bb = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    except (ValueError, IndexError):
        return a
    t = max(0.0, min(1.0, t))
    return "%02X%02X%02X" % (
        int(round(ar + (br - ar) * t)),
        int(round(ag + (bg - ag) * t)),
        int(round(ab + (bb - ab) * t)),
    )


def _studio_ramp(el):
    key = str(el.get("ramp") or "accent")
    if key in _STUDIO_RAMPS:
        return _STUDIO_RAMPS[key]
    accent = _clean_hex(el.get("accent")) or "1D4E89"
    return [
        _mix_hex(accent, "FFFFFF", 0.88), _mix_hex(accent, "FFFFFF", 0.66),
        _mix_hex(accent, "FFFFFF", 0.38), accent,
        _mix_hex(accent, "0B1220", 0.22), _mix_hex(accent, "0B1220", 0.44),
    ]


def _ramp_at(stops, t):
    t = max(0.0, min(1.0, t))
    n = len(stops) - 1
    i = min(n - 1, int(t * n))
    return _mix_hex(stops[i], stops[i + 1], t * n - i)


def _series_color(el, i, n):
    """Categorical colour. Mirrors seriesColor() in hanns_core.js: the palest
    end of a ramp is fine as a map fill but illegible as a bar or a label,
    so series sample the upper part only."""
    stops = _studio_ramp(el)
    t = (i / (n - 1)) if n > 1 else 0.5
    return _ramp_at(stops, 0.30 + t * 0.66)


def _readable_on(hex6):
    try:
        r, g, b = int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16)
    except (ValueError, IndexError):
        return "FFFFFF"
    return "10202E" if (0.299 * r + 0.587 * g + 0.114 * b) > 150 else "FFFFFF"


def _studio_rows(el):
    out = []
    for r in (el.get("rows") or []):
        if not isinstance(r, dict):
            continue
        out.append({
            "label": "" if r.get("label") is None else str(r.get("label")),
            "value": _num(r.get("value"), 0.0),
            "value2": None if r.get("value2") in (None, "") else _num(r.get("value2"), 0.0),
            "color": _clean_hex(r.get("color")),
            "note": "" if r.get("note") is None else str(r.get("note")),
        })
    return out


def _studio_fmt(el, v):
    dec = int(_num(el.get("decimals"), 0))
    dec = 0 if dec < 0 else (4 if dec > 4 else dec)
    try:
        body = f"{float(v):,.{dec}f}"
    except (TypeError, ValueError):
        body = str(v)
    return f"{el.get('valuePrefix') or ''}{body}{el.get('valueSuffix') or ''}"


def _txt_box(slide, x, y, w, h, text, *, size=14, bold=False, color="0F172A",
             align=PP_ALIGN.LEFT, anchor=None, italic=False):
    """A plain text box in slide pixel coordinates."""
    box = slide.shapes.add_textbox(_px_to_emu(x), _px_to_emu(y),
                                   _px_to_emu(max(8, w)), _px_to_emu(max(8, h)))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(max(6, int(round(size / 1.333))))
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = _rgb("#" + (_clean_hex(color) or "0F172A"), "0F172A")
    return box


def _rect(slide, x, y, w, h, fill, *, radius=0, line=None, line_w=0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        _px_to_emu(x), _px_to_emu(y), _px_to_emu(max(2, w)), _px_to_emu(max(2, h)))
    if radius:
        _set_rounded_radius(shp, radius, min(w, h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        shp.fill.background()
    if line and line_w:
        shp.line.color.rgb = RGBColor.from_string(line)
        shp.line.width = Pt(max(0.5, line_w / 1.333))
    else:
        shp.line.fill.background()
    return shp


def _studio_frame(el):
    """Element box in slide pixels, plus a cursor for stacked content."""
    x, y = _num(el.get("x"), 0), _num(el.get("y"), 0)
    w, h = max(40, _num(el.get("w"), 300)), max(30, _num(el.get("h"), 200))
    return x, y, w, h


def _studio_title(slide, el, x, y, w):
    """Draw the object title if it has one; return the height consumed."""
    title = str(el.get("title") or "").strip()
    if not title:
        return 0
    dark = bool(el.get("dark"))
    _txt_box(slide, x, y, w, 26, title, size=20, bold=True,
             color="F1F5F9" if dark else "0F172A")
    return 30


def _std_stat_block(slide, el):
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el) or [{"label": "", "value": _num(el.get("level"), 0),
                                 "value2": None, "note": "", "color": None}]
    r = rows[0]
    accent = _clean_hex(el.get("accent")) or "E8482B"
    dark = bool(el.get("dark"))
    align = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
        str(el.get("textAlign") or "left"), PP_ALIGN.LEFT)
    cy = y
    if el.get("title"):
        _txt_box(slide, x, cy, w, 18, str(el["title"]).upper(), size=13, bold=True,
                 color=accent, align=align)
        cy += 22
    vh = min(h * 0.52, 96)
    _txt_box(slide, x, cy, w, vh, _studio_fmt(el, r["value"]),
             size=max(28, vh * 0.92), bold=True, color=accent, align=align)
    cy += vh + 4
    if r["label"]:
        _txt_box(slide, x, cy, w, 40, r["label"], size=18, bold=True,
                 color="E6EDF5" if dark else "0F172A", align=align)
        cy += 42
    if r["value2"] is not None:
        up = r["value2"] >= 0
        _txt_box(slide, x, cy, w, 20,
                 ("\u25B2 " if up else "\u25BC ") + _studio_fmt(el, abs(r["value2"])),
                 size=14, bold=True, color="15803D" if up else "B91C1C", align=align)
        cy += 22
    if r["note"]:
        _txt_box(slide, x, cy, w, 24, r["note"], size=13,
                 color="94A3B8" if dark else "5B7183", align=align)
    return None


def _std_kpi_grid(slide, el):
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el)
    if not rows:
        return None
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    cols = max(1, min(6, int(_num(el.get("cols"), min(3, len(rows))))))
    import math as _math
    rws = _math.ceil(len(rows) / cols)
    gap = 10
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rws - 1)) / max(1, rws)
    style = str(el.get("tileStyle") or "soft")
    for i, r in enumerate(rows):
        cx = x + (i % cols) * (cw + gap)
        cyy = y + (i // cols) * (ch + gap)
        tint = r["color"] or _series_color(el, i, len(rows))
        if style == "solid":
            _rect(slide, cx, cyy, cw, ch, tint, radius=14)
            ink = _readable_on(tint)
        else:
            _rect(slide, cx, cyy, cw, ch, _mix_hex(tint, "FFFFFF", 0.85), radius=14,
                  line=tint, line_w=1.2)
            ink = tint
        _txt_box(slide, cx + 12, cyy + 10, cw - 24, ch * 0.5,
                 _studio_fmt(el, r["value"]), size=min(46, ch * 0.42), bold=True, color=ink)
        if r["label"]:
            _txt_box(slide, cx + 12, cyy + ch * 0.55, cw - 24, ch * 0.24, r["label"],
                     size=min(15, ch * 0.16), bold=True,
                     color=ink if style == "solid" else "0F172A")
        if r["note"]:
            _txt_box(slide, cx + 12, cyy + ch * 0.76, cw - 24, ch * 0.2, r["note"],
                     size=min(12, ch * 0.13),
                     color=ink if style == "solid" else "5B7183")
    return None


def _std_bar_rows(slide, el, ranked):
    """Ranked bars and bullet bars: same row geometry, different trimmings."""
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el)
    if not rows:
        return None
    if ranked:
        srt = str(el.get("sort") or "desc")
        if srt == "desc":
            rows.sort(key=lambda r: -r["value"])
        elif srt == "asc":
            rows.sort(key=lambda r: r["value"])
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    hi = _num(el.get("max"), 0) or max([r["value"] for r in rows] +
                                       [(r["value2"] or 0) for r in rows] + [1])
    n = len(rows)
    gap = 7
    rh = max(12, (h - gap * (n - 1)) / n)
    pos_w = 26 if (ranked and el.get("showNumbers") is not False) else 0
    lab_w = max(70, w * 0.24)
    val_w = 66 if el.get("showValues") is not False else 0
    track_x = x + pos_w + lab_w + 8
    track_w = max(30, w - pos_w - lab_w - val_w - 16)
    dark = bool(el.get("dark"))
    for i, r in enumerate(rows):
        ry = y + i * (rh + gap)
        c = r["color"] or _series_color(el, (n - 1 - i) if ranked else i, n)
        if pos_w:
            _rect(slide, x, ry + rh * 0.1, 22, rh * 0.8,
                  "1F293722" if dark else "E2E8F0", radius=6)
            _txt_box(slide, x, ry + rh * 0.22, 22, rh * 0.6, str(i + 1), size=min(13, rh * 0.5),
                     bold=True, color="0F172A", align=PP_ALIGN.CENTER)
        _txt_box(slide, x + pos_w + 4, ry + rh * 0.15, lab_w, rh * 0.7, r["label"],
                 size=min(16, rh * 0.56), bold=True,
                 color="E6EDF5" if dark else "0F172A")
        _rect(slide, track_x, ry + rh * 0.22, track_w, rh * 0.56,
              "334155" if dark else "EDF1F6", radius=rh * 0.28)
        fw = max(2, track_w * max(0.0, min(1.0, r["value"] / hi)))
        _rect(slide, track_x, ry + rh * 0.22, fw, rh * 0.56, c, radius=rh * 0.28)
        if r["value2"] is not None:
            tx = track_x + track_w * max(0.0, min(1.0, r["value2"] / hi))
            _rect(slide, tx - 1.5, ry + rh * 0.1, 3, rh * 0.8,
                  "E2E8F0" if dark else "0F172A")
        if val_w:
            _txt_box(slide, x + w - val_w, ry + rh * 0.15, val_w, rh * 0.7,
                     _studio_fmt(el, r["value"]), size=min(16, rh * 0.56), bold=True,
                     color="E6EDF5" if dark else "0F172A", align=PP_ALIGN.RIGHT)
    return None


def _std_process_steps(slide, el):
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el)
    if not rows:
        return None
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    n = len(rows)
    vertical = str(el.get("orient") or "horizontal") == "vertical"
    chevron = str(el.get("stepStyle") or "chevron") == "chevron" and not vertical
    gap = 6
    if vertical:
        sh = (h - gap * (n - 1)) / n
        sw = w
    else:
        sw = (w - gap * (n - 1)) / n
        sh = h
    for i, r in enumerate(rows):
        sx = x + (0 if vertical else i * (sw + gap))
        sy = y + (i * (sh + gap) if vertical else 0)
        tint = r["color"] or _series_color(el, i, n)
        if chevron:
            shp = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON, _px_to_emu(sx), _px_to_emu(sy),
                _px_to_emu(max(4, sw)), _px_to_emu(max(4, sh)))
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor.from_string(tint)
            shp.line.fill.background()
        else:
            _rect(slide, sx, sy, sw, sh, _mix_hex(tint, "FFFFFF", 0.84),
                  radius=12, line=tint, line_w=1.2)
        ink = _readable_on(tint) if chevron else "0F172A"
        pad = sh * 0.22 if chevron else 12
        label = (f"{i + 1}. " if el.get("showNumbers") is not False else "") + r["label"]
        _txt_box(slide, sx + pad, sy + sh * 0.22, sw - pad * 2, sh * 0.32, label,
                 size=min(17, sh * 0.24), bold=True, color=ink,
                 align=PP_ALIGN.CENTER if chevron else PP_ALIGN.LEFT)
        if r["note"]:
            _txt_box(slide, sx + pad, sy + sh * 0.55, sw - pad * 2, sh * 0.3, r["note"],
                     size=min(13, sh * 0.17), color=ink,
                     align=PP_ALIGN.CENTER if chevron else PP_ALIGN.LEFT)
    return None


def _std_timeline(slide, el):
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el)
    if not rows:
        return None
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    n = len(rows)
    dark = bool(el.get("dark"))
    axis_y = y + h / 2
    _rect(slide, x, axis_y - 1.5, w, 3, "334155" if dark else "D8E0E8")
    step = w / max(1, n)
    for i, r in enumerate(rows):
        cx = x + step * i + step / 2
        tint = r["color"] or _series_color(el, i, n)
        slide.shapes.add_shape(MSO_SHAPE.OVAL, _px_to_emu(cx - 7), _px_to_emu(axis_y - 7),
                               _px_to_emu(14), _px_to_emu(14)).fill.solid()
        dot = slide.shapes[-1]
        dot.fill.fore_color.rgb = RGBColor.from_string(tint)
        dot.line.fill.background()
        above = (i % 2 == 0) or el.get("alternate") is False
        ty = (axis_y - h * 0.44) if above else (axis_y + 16)
        _txt_box(slide, cx - step / 2, ty, step, 22, r["label"], size=17, bold=True,
                 color=tint, align=PP_ALIGN.CENTER)
        if r["note"]:
            _txt_box(slide, cx - step / 2, ty + 22, step, 34, r["note"], size=12,
                     color="94A3B8" if dark else "5B7183", align=PP_ALIGN.CENTER)
    return None


def _arc(slide, x, y, d, start_deg, end_deg, fill, thickness=0.26):
    """A BLOCK_ARC is a ring with a real sweep: adjustments are
    (start angle, end angle, thickness as a fraction of the radius)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, _px_to_emu(x), _px_to_emu(y),
                                 _px_to_emu(d), _px_to_emu(d))
    try:
        shp.adjustments[0] = float(start_deg)
        shp.adjustments[1] = float(end_deg)
        shp.adjustments[2] = float(thickness)
    except (IndexError, ValueError):
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    shp.line.fill.background()
    return shp


def _std_ring_grid(slide, el):
    """Percentage donuts. PowerPoint can draw a genuine partial ring, so the
    sweep encodes the value rather than the number doing all the work."""
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el)
    if not rows:
        return None
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    cols = max(1, min(6, int(_num(el.get("cols"), min(4, len(rows))))))
    gap = 10
    cw = (w - gap * (cols - 1)) / cols
    hi = _num(el.get("max"), 100) or 100
    dark = bool(el.get("dark"))
    thickness = max(0.08, min(0.45, _num(el.get("ringW"), 13) / 50.0))
    for i, r in enumerate(rows):
        cx = x + (i % cols) * (cw + gap)
        d = max(24, min(cw, h * 0.7))
        ox = cx + (cw - d) / 2
        tint = r["color"] or _series_color(el, i, len(rows))
        pct = max(0.0, min(1.0, (r["value"] / hi) if hi else 0.0))
        # Track, then the sweep clockwise from 12 o'clock.
        _arc(slide, ox, y, d, 0, 359.9, "334155" if dark else "E6EBF1", thickness)
        if pct > 0.001:
            _arc(slide, ox, y, d, 270.0, (270.0 + 359.9 * pct) % 360.0, tint, thickness)
        _txt_box(slide, ox, y + d * 0.38, d, d * 0.26, _studio_fmt(el, r["value"]),
                 size=max(10, min(26, d * 0.24)), bold=True,
                 color="FFFFFF" if dark else "0F172A", align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        if r["label"]:
            _txt_box(slide, cx, y + d + 6, cw, 20, r["label"], size=13, bold=True,
                     color="E6EDF5" if dark else "0F172A", align=PP_ALIGN.CENTER)
    return None


def _std_waffle(slide, el):
    x, y, w, h = _studio_frame(el)
    rows = _studio_rows(el)
    cols = max(2, min(20, int(_num(el.get("cols"), 10))))
    total = cols * cols
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    side = min(w, h)
    cell = side / cols
    pad = cell * 0.14
    cells = []
    if len(rows) <= 1:
        pct = max(0.0, min(100.0, rows[0]["value"] if rows else _num(el.get("level"), 60)))
        on = int(round(pct / 100 * total))
        c = (rows[0]["color"] if rows else None) or _clean_hex(el.get("accent")) or "0284C7"
        cells = [c if i < on else None for i in range(total)]
    else:
        ssum = sum(max(0.0, r["value"]) for r in rows) or 1
        acc = []
        for i, r in enumerate(rows):
            c = r["color"] or _series_color(el, i, len(rows))
            acc += [c] * int(round(max(0.0, r["value"]) / ssum * total))
        cells = [acc[i] if i < len(acc) else None for i in range(total)]
    dark = bool(el.get("dark"))
    for i, c in enumerate(cells):
        cx = x + (i % cols) * cell
        cy = y + (i // cols) * cell
        _rect(slide, cx, cy, cell - pad, cell - pad,
              c or ("2B3648" if dark else "E4E9EF"), radius=max(1, cell * 0.18))
    return None


def _std_heat_grid(slide, el):
    """A labelled matrix maps cleanly onto a real PowerPoint table."""
    rows = _studio_rows(el)
    if not rows:
        return None
    cells = [[_num(v, None) for v in re.split(r"[,;\s]+", r["note"]) if v != ""] for r in rows]
    ncols = max([len(c) for c in cells] + [1])
    col_labels = [s.strip() for s in str(el.get("colLabels") or "").split(",")]
    flat = [v for row in cells for v in row if v is not None]
    lo = _num(el.get("scaleMin"), min(flat) if flat else 0)
    hi = _num(el.get("scaleMax"), max(flat) if flat else 1)
    if hi == lo:
        hi = lo + 1
    stops = _studio_ramp(el)
    x, y, w, h = _studio_frame(el)
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    shape = slide.shapes.add_table(len(rows) + 1, ncols + 1, _px_to_emu(x), _px_to_emu(y),
                                   _px_to_emu(w), _px_to_emu(h))
    tbl = shape.table
    tbl.cell(0, 0).text = ""
    for c in range(ncols):
        tbl.cell(0, c + 1).text = col_labels[c] if c < len(col_labels) else ""
    for ri, r in enumerate(rows):
        tbl.cell(ri + 1, 0).text = r["label"]
        for c in range(ncols):
            cell = tbl.cell(ri + 1, c + 1)
            v = cells[ri][c] if c < len(cells[ri]) else None
            if v is None:
                cell.text = ""
                continue
            bg = _ramp_at(stops, (v - lo) / (hi - lo))
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(bg)
            cell.text = "" if el.get("showValues") is False else _studio_fmt(el, v)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor.from_string(_readable_on(bg))
    return shape


def _std_quote_card(slide, el):
    x, y, w, h = _studio_frame(el)
    accent = _clean_hex(el.get("accent")) or "7C3AED"
    dark = bool(el.get("dark"))
    style = str(el.get("quoteStyle") or "bar")
    if style == "bar":
        _rect(slide, x, y, 6, h, accent)
        x += 22
        w -= 22
    elif style == "card":
        _rect(slide, x, y, w, h, None, radius=18, line=accent, line_w=2)
        x += 18
        w -= 36
    quote = str(el.get("quote") or "")
    _txt_box(slide, x, y + h * 0.14, w, h * 0.5, "\u201C" + quote + "\u201D",
             size=_num(el.get("quoteSize"), 34), bold=True,
             color="E6EDF5" if dark else "0F172A")
    cy = y + h * 0.66
    if el.get("attribution"):
        _txt_box(slide, x, cy, w, 24, str(el["attribution"]), size=18, bold=True, color=accent)
        cy += 26
    if el.get("role"):
        _txt_box(slide, x, cy, w, 22, str(el["role"]), size=14,
                 color="94A3B8" if dark else "5B7183")
    return None


def _std_data_card(slide, el, kind_label):
    """Fallback for the geometric objects. A choropleth or a Sankey has no
    honest native PowerPoint equivalent, so the deck gets a titled card with
    the underlying figures listed — the numbers survive, editable."""
    x, y, w, h = _studio_frame(el)
    accent = _clean_hex(el.get("accent")) or "94A3B8"
    _rect(slide, x, y, w, h, "F6F9FC", radius=18, line=accent, line_w=1.5)
    title = str(el.get("title") or el.get("label") or kind_label)
    _txt_box(slide, x + 16, y + 12, w - 32, 26, title, size=19, bold=True, color="0F172A")
    rows = _studio_rows(el)
    if not rows:
        return None
    lines = [f"{r['label']}   {_studio_fmt(el, r['value'])}" for r in rows[:14]]
    if len(rows) > 14:
        lines.append(f"… and {len(rows) - 14} more")
    body = slide.shapes.add_textbox(_px_to_emu(x + 16), _px_to_emu(y + 44),
                                    _px_to_emu(max(20, w - 32)), _px_to_emu(max(20, h - 56)))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor.from_string("334155")
    return None


# ── mechanical objects ────────────────────────────────────────────────
# These animate in the browser. A slide is a still frame, so each exports
# at the state it is set to: the clock at its time, the gauge at its value,
# the hourglass part-run. Nothing pretends to move.

def _std_clock_face(slide, el):
    import math
    x, y, w, h = _studio_frame(el)
    d = min(w, h)
    cx, cy = x + w / 2, y + d / 2
    r = d / 2
    accent = _clean_hex(el.get("accent")) or "1D4E89"
    dark = bool(el.get("dark"))
    face = _clean_hex(el.get("faceColor")) or ("0F172A" if dark else "FFFFFF")
    ink = _clean_hex(el.get("handColor")) or ("E6EDF5" if dark else "0F172A")
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, _px_to_emu(cx - r), _px_to_emu(cy - r),
                                  _px_to_emu(d), _px_to_emu(d))
    ring.fill.solid()
    ring.fill.fore_color.rgb = RGBColor.from_string(face)
    ring.line.color.rgb = RGBColor.from_string(accent)
    ring.line.width = Pt(max(1.0, _num(el.get("bezelW"), 7) / 1.333))

    mode = str(el.get("clockMode") or "live")
    if mode == "fixed":
        hh, mm = int(_num(el.get("hour"), 10)) % 12, int(_num(el.get("minute"), 10)) % 60
    else:
        # A static export of a "live" clock is the moment it was exported.
        from datetime import datetime
        now = datetime.now()
        hh, mm = now.hour % 12, now.minute
    if str(el.get("faceStyle") or "ticks") in ("numbers", "roman"):
        RN = ["XII", "I", "II", "III", "IV", "V", "VI", "VII",
              "VIII", "IX", "X", "XI"]
        roman = str(el.get("faceStyle")) == "roman"
        for i in range(12):
            a = math.radians(i * 30 - 90)
            lx, ly = cx + math.cos(a) * r * 0.72, cy + math.sin(a) * r * 0.72
            _txt_box(slide, lx - 22, ly - 12, 44, 24,
                     RN[i] if roman else ("12" if i == 0 else str(i)),
                     size=max(9, r * 0.19), bold=True, color=ink, align=PP_ALIGN.CENTER)
    for ang, length, width in ((( hh + mm / 60.0) * 30 - 90, r * 0.52, 6.0),
                               (mm * 6 - 90, r * 0.76, 4.0)):
        a = math.radians(ang)
        ln = slide.shapes.add_connector(
            1, _px_to_emu(cx), _px_to_emu(cy),
            _px_to_emu(cx + math.cos(a) * length), _px_to_emu(cy + math.sin(a) * length))
        ln.line.color.rgb = RGBColor.from_string(ink)
        ln.line.width = Pt(width)
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, _px_to_emu(cx - 6), _px_to_emu(cy - 6),
                                 _px_to_emu(12), _px_to_emu(12))
    hub.fill.solid()
    hub.fill.fore_color.rgb = RGBColor.from_string(accent)
    hub.line.fill.background()
    return ring


def _std_gears(slide, el):
    """Meshing wheels as real PPTX gear autoshapes, laid out left to right
    with each wheel sized by its tooth count."""
    rows = _studio_rows(el) or [{"label": "", "value": 18, "color": None, "note": "", "value2": None},
                                {"label": "", "value": 12, "color": None, "note": "", "value2": None}]
    x, y, w, h = _studio_frame(el)
    top = _studio_title(slide, el, x, y, w)
    y += top
    h -= top
    teeth = [max(6, min(40, int(abs(_num(r["value"], 12))))) for r in rows]
    radii = [t * 3.4 for t in teeth]
    centres, run = [], radii[0]
    centres.append(run)
    for i in range(1, len(radii)):
        run += (radii[i - 1] + radii[i]) * 0.90
        centres.append(run)
    span = centres[-1] + radii[-1]
    label_room = 22 if any(r["label"] for r in rows) else 0
    k = min((w * 0.94) / max(1, span), (h * 0.9 - label_room) / max(1, max(radii) * 2))
    ox = x + w / 2 - (span * k) / 2
    cy = y + (h - label_room) / 2
    # python-pptx exposes GEAR_6 and GEAR_9; pick whichever is closer.
    for i, r in enumerate(rows):
        R = radii[i] * k
        cx = ox + centres[i] * k
        col = r["color"] or _series_color(el, i, len(rows))
        shape = MSO_SHAPE.GEAR_9 if teeth[i] >= 12 else MSO_SHAPE.GEAR_6
        g = slide.shapes.add_shape(shape, _px_to_emu(cx - R), _px_to_emu(cy - R),
                                   _px_to_emu(R * 2), _px_to_emu(R * 2))
        g.fill.solid()
        g.fill.fore_color.rgb = RGBColor.from_string(col)
        g.line.fill.background()
        if r["label"]:
            _txt_box(slide, cx - R, cy + R + 2, R * 2, 20, r["label"], size=13, bold=True,
                     color="E6EDF5" if el.get("dark") else "0F172A", align=PP_ALIGN.CENTER)
    return None


def _std_battery(slide, el):
    x, y, w, h = _studio_frame(el)
    r = (_studio_rows(el) or [{"label": "", "value": _num(el.get("level"), 72)}])[0]
    pct = max(0.0, min(100.0, r["value"]))
    auto = "DC2626" if pct < 15 else ("EA580C" if pct < 35 else "16A34A")
    c = r.get("color") or _clean_hex(el.get("accent")) or auto
    frame = "E2E8F0" if el.get("dark") else "1F2937"
    vertical = str(el.get("orient") or "horizontal") == "vertical"
    if vertical:
        bw, bh = min(w * 0.62, h * 0.42), h * 0.88
        bx, by = x + (w - bw) / 2, y + h * 0.1
        _rect(slide, bx + bw * 0.32, by - h * 0.05, bw * 0.36, h * 0.05, frame, radius=4)
        _rect(slide, bx, by, bw, bh, None, radius=16, line=frame, line_w=5)
        fh = (bh - 12) * pct / 100
        _rect(slide, bx + 6, by + bh - 6 - fh, bw - 12, fh, c, radius=10)
    else:
        bw, bh = w * 0.88, min(h * 0.72, w * 0.5)
        bx, by = x + w * 0.03, y + (h - bh) / 2
        _rect(slide, bx + bw + 2, by + bh * 0.32, w * 0.05, bh * 0.36, frame, radius=4)
        _rect(slide, bx, by, bw, bh, None, radius=16, line=frame, line_w=5)
        _rect(slide, bx + 6, by + 6, (bw - 12) * pct / 100, bh - 12, c, radius=10)
    if el.get("showValues") is not False:
        _txt_box(slide, x, y + h * 0.38, w, h * 0.26, _studio_fmt(el, pct),
                 size=min(38, h * 0.3), bold=True,
                 color="FFFFFF" if pct > 45 else ("E6EDF5" if el.get("dark") else "0F172A"),
                 align=PP_ALIGN.CENTER)
    if r["label"]:
        _txt_box(slide, x, y + h - 18, w, 18, r["label"], size=12, bold=True,
                 color="94A3B8" if el.get("dark") else "5B7183", align=PP_ALIGN.CENTER)
    return None


def _std_thermometer(slide, el):
    x, y, w, h = _studio_frame(el)
    r = (_studio_rows(el) or [{"label": "", "value": _num(el.get("level"), 28)}])[0]
    lo, hi = _num(el.get("scaleMin"), 0), _num(el.get("scaleMax"), 50)
    t = max(0.0, min(1.0, (r["value"] - lo) / max(1e-9, hi - lo)))
    c = r.get("color") or _clean_hex(el.get("accent")) or "DC2626"
    frame = "CBD5E1" if el.get("dark") else "334155"
    tube_w = min(w * 0.22, 34)
    cx = x + w * 0.4
    top_y, bulb_r = y + h * 0.06, tube_w * 0.95
    bot_y = y + h - bulb_r * 2.1
    _rect(slide, cx - tube_w / 2, top_y, tube_w, bot_y - top_y + bulb_r,
          "FFFFFF", radius=tube_w / 2, line=frame, line_w=3)
    fh = (bot_y - top_y) * t
    _rect(slide, cx - tube_w / 2 + 4, bot_y - fh, tube_w - 8, fh + bulb_r, c,
          radius=(tube_w - 8) / 2)
    bulb = slide.shapes.add_shape(MSO_SHAPE.OVAL, _px_to_emu(cx - bulb_r),
                                  _px_to_emu(bot_y + bulb_r * 0.1),
                                  _px_to_emu(bulb_r * 2), _px_to_emu(bulb_r * 2))
    bulb.fill.solid()
    bulb.fill.fore_color.rgb = RGBColor.from_string(c)
    bulb.line.color.rgb = RGBColor.from_string(frame)
    bulb.line.width = Pt(2)
    ticks = max(2, min(11, int(_num(el.get("legendTicks"), 6))))
    for i in range(ticks):
        f = i / (ticks - 1)
        ty = bot_y - (bot_y - top_y) * f
        _txt_box(slide, cx + tube_w, ty - 9, w * 0.5, 18, _studio_fmt(el, lo + (hi - lo) * f),
                 size=11, bold=True, color="94A3B8" if el.get("dark") else "5B7183")
    if r["label"]:
        _txt_box(slide, x, y + h - 16, w, 16, r["label"], size=12, bold=True,
                 color="94A3B8" if el.get("dark") else "5B7183", align=PP_ALIGN.CENTER)
    return None


def _std_speedometer(slide, el):
    """The dial is a BLOCK_ARC track plus a coloured sweep, with the reading
    below the hub — the same layout the browser draws."""
    x, y, w, h = _studio_frame(el)
    r = (_studio_rows(el) or [{"label": "", "value": _num(el.get("level"), 68)}])[0]
    lo, hi = _num(el.get("scaleMin"), 0), _num(el.get("scaleMax"), 100)
    t = max(0.0, min(1.0, (r["value"] - lo) / max(1e-9, hi - lo)))
    accent = _clean_hex(el.get("accent")) or "0D9488"
    d = min(w, h * 1.5)
    cx, cy = x + w / 2, y + h * 0.52
    START, SPAN = 160.0, 220.0
    _arc(slide, cx - d / 2, cy - d / 2, d, START, START + SPAN,
         "334155" if el.get("dark") else "E6EBF1", 0.17)
    if t > 0.004:
        _arc(slide, cx - d / 2, cy - d / 2, d, START, (START + SPAN * t) % 360.0, accent, 0.17)
    if el.get("showValues") is not False:
        _txt_box(slide, x, cy + d * 0.10, w, d * 0.22, _studio_fmt(el, r["value"]),
                 size=min(40, d * 0.19), bold=True,
                 color="F1F5F9" if el.get("dark") else "0F172A", align=PP_ALIGN.CENTER)
    if r["label"]:
        _txt_box(slide, x, cy + d * 0.33, w, 20, r["label"], size=13, bold=True,
                 color="94A3B8" if el.get("dark") else "5B7183", align=PP_ALIGN.CENTER)
    return None


def _std_hourglass(slide, el):
    """Drawn mid-run: sand in both bulbs, which is what an hourglass on a
    still slide should look like."""
    x, y, w, h = _studio_frame(el)
    sand = _clean_hex(el.get("sandColor")) or _clean_hex(el.get("accent")) or "C2861A"
    frame = _clean_hex(el.get("frameColor")) or ("E2E8F0" if el.get("dark") else "3F2F1C")
    gw = min(w * 0.8, h * 0.62)
    gx, gy = x + (w - gw) / 2, y + h * 0.06
    gh = h * 0.88
    cap = max(8, gh * 0.06)
    _rect(slide, gx - gw * 0.06, gy, gw * 1.12, cap, frame, radius=cap / 2)
    _rect(slide, gx - gw * 0.06, gy + gh - cap, gw * 1.12, cap, frame, radius=cap / 2)
    up = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, _px_to_emu(gx),
                                _px_to_emu(gy + cap), _px_to_emu(gw), _px_to_emu((gh - cap * 2) / 2))
    up.rotation = 180
    dn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, _px_to_emu(gx),
                                _px_to_emu(gy + cap + (gh - cap * 2) / 2),
                                _px_to_emu(gw), _px_to_emu((gh - cap * 2) / 2))
    for shp, frac in ((up, 0.55), (dn, 0.45)):
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(sand)
        shp.line.color.rgb = RGBColor.from_string(frame)
        shp.line.width = Pt(2)
    if el.get("title"):
        _txt_box(slide, x, y + h - 18, w, 18, str(el["title"]), size=12, bold=True,
                 color="94A3B8" if el.get("dark") else "5B7183", align=PP_ALIGN.CENTER)
    return None


_STUDIO_NATIVE = {
    "stat_block":     _std_stat_block,
    "kpi_grid":       _std_kpi_grid,
    "rank_bars":      lambda s, e: _std_bar_rows(s, e, True),
    "bullet_bars":    lambda s, e: _std_bar_rows(s, e, False),
    "process_steps":  _std_process_steps,
    "timeline_track": _std_timeline,
    "ring_grid":      _std_ring_grid,
    "waffle":         _std_waffle,
    "heat_grid":      _std_heat_grid,
    "quote_card":     _std_quote_card,
    "clock_face":     _std_clock_face,
    "gears":          _std_gears,
    "charge_meter":   _std_battery,
    "temp_gauge":     _std_thermometer,
    "speedometer":    _std_speedometer,
    "sand_timer":     _std_hourglass,
}
_STUDIO_CARD_LABEL = {
    "choropleth": "Region map", "gradient_legend": "Colour scale",
    "slope_chart": "Slope chart", "matrix_2x2": "2x2 matrix",
    "venn": "Venn diagram", "pyramid_tiers": "Pyramid", "sankey_flow": "Flow",
}


def _add_studio_object(slide, el):
    kind = str(el.get("objectType") or "")
    fn = _STUDIO_NATIVE.get(kind)
    try:
        if fn:
            return fn(slide, el)
        return _std_data_card(slide, el, _STUDIO_CARD_LABEL.get(kind, "Infographic"))
    except Exception:
        # Never let one object take the whole export down.
        return _add_rich_placeholder(slide, el, _STUDIO_CARD_LABEL.get(kind, "Infographic"))


def _render_element(slide, el, fx: _AnimFx):
    """Render one element and register its entrance animation."""
    etype = str(el.get("type", "")).lower()
    shp = None
    if etype == "text":
        shp = _add_text(slide, el)
    elif etype == "image":
        shp = _add_image(slide, el)
    elif etype in ("rect", "rectangle"):
        shp = _add_shape(slide, el, "rect")
    elif etype == "ellipse":
        shp = _add_shape(slide, el, "ellipse")
    elif etype == "line":
        shp = _add_line(slide, el)
    elif etype == "chart":
        shp = _add_chart(slide, el)
    elif etype == "table":
        shp = _add_table(slide, el)
    elif etype == "link":
        shp = _add_link(slide, el)
    elif etype == "video":
        shp = _add_video(slide, el)
    elif etype == "map":
        shp = _add_rich_placeholder(slide, el, "Map")
    elif etype == "object":
        if str(el.get("objectType") or "") in STUDIO_KINDS:
            # Studio objects are rebuilt from native shapes so they stay
            # editable in PowerPoint rather than exporting as a flat card.
            shp = _add_studio_object(slide, el)
        else:
            shp = _add_rich_placeholder(slide, el, "Infographic")
    elif etype == "creative_shape":
        shp = _add_shape(slide, {**el, "type": "ellipse", "radius": 0}, "ellipse")
    elif etype == "group":
        # Children are stored in slide coordinates offset by the group box.
        gx, gy = _num(el.get("x", 0)), _num(el.get("y", 0))
        for child in (el.get("children") or []):
            if not isinstance(child, dict):
                continue
            shifted = dict(child)
            shifted["x"] = _num(child.get("x", 0)) + gx
            shifted["y"] = _num(child.get("y", 0)) + gy
            # Group-level animation applies to each child.
            if el.get("anim") and el.get("anim") != "none":
                shifted.setdefault("anim", el.get("anim"))
                shifted.setdefault("animDelay", el.get("animDelay", 0))
            _render_element(slide, shifted, fx)
        return
    else:
        if el.get("text"):
            shp = _add_text(slide, {**el, "type": "text"})

    if shp is not None:
        fx.register(shp, el)


# ───────────────────────── public API ─────────────────────────

def export_deck_to_pptx(deck) -> io.BytesIO:
    """Build a .pptx for ``deck`` and return it as an in-memory BytesIO."""
    if Presentation is None:
        raise ValueError(
            "python-pptx is not installed. Add python-pptx to requirements.txt and rebuild."
        )

    prs = Presentation()
    prs.slide_width = Emu(_px_to_emu(DESIGN_W))
    prs.slide_height = Emu(_px_to_emu(DESIGN_H))
    blank = prs.slide_layouts[6]  # fully blank layout

    slides = list(deck.slides.all())
    if not slides:
        prs.slides.add_slide(blank)

    for srow in slides:
        data = srow.as_dict() if hasattr(srow, "as_dict") else dict(srow.data or {})
        slide = prs.slides.add_slide(blank)
        _set_slide_bg(slide, data)

        fx = _AnimFx()
        # PPTX z-order is document order, and Hanns paints elements in array
        # order — so exporting in the SAME order keeps stacking identical to
        # the live stage.
        for el in (data.get("els") or []):
            try:
                _render_element(slide, el, fx)
            except Exception:
                # One bad element must never abort the whole export.
                continue

        fx.apply(slide)
        _set_notes(slide, data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def export_filename(deck) -> str:
    base = re.sub(r"[^A-Za-z0-9 _-]+", "", (deck.title or "deck")).strip() or "deck"
    return f"{base}.pptx"